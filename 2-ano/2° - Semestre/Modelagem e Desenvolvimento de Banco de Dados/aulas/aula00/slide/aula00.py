from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    return slide

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions
    return slide

# 1 — Abertura
add_slide(
    "Introdução aos Bancos de Dados",
    "Fundamentos de Banco de Dados\n\n"
    "Curso Técnico em Informática\n"
    "Aula 1 — 50 minutos"
)

# 2 — Problema inicial
add_activity_slide(
    "Para começar: onde guardaríamos esses dados?",
    "Imagine que precisamos guardar os dados de todos os alunos da escola.\n\n"
    "• 10 alunos: como você faria?\n"
    "• 100 alunos: ainda funcionaria?\n"
    "• 10.000 alunos: e agora?\n"
    "• Como encontrar um aluno específico?\n"
    "• Como alterar ou excluir seus dados?"
)

# 3 — Dados e informação
add_slide(
    "Dados × Informação",
    "DADO\n"
    "Um valor ou fato que pode ser registrado.\n\n"
    "Exemplo:\n"
    "João | 2005 | Informática | 12345\n\n"
    "INFORMAÇÃO\n"
    "Quando os dados possuem contexto e significado:\n"
    "Nome: João\n"
    "Ano de nascimento: 2005\n"
    "Curso: Informática\n"
    "Matrícula: 12345"
)

# 4 — O que é BD
add_slide(
    "O que é um Banco de Dados?",
    "É uma coleção organizada de dados relacionados,\n"
    "com significado e propósito.\n\n"
    "Um banco de dados não é simplesmente um conjunto\n"
    "aleatório de valores.\n\n"
    "Dados + organização + significado + propósito"
)

# 5 — Exemplo de aluno
add_slide(
    "Exemplo: dados de um aluno",
    "ALUNO\n\n"
    "• Nome\n"
    "• Matrícula\n"
    "• RG\n"
    "• Curso\n"
    "• Ano de ingresso\n"
    "• Endereço\n\n"
    "Esses dados representam uma parte da realidade\n"
    "que interessa ao sistema."
)

# 6 — Minimundo
add_slide(
    "Minimundo",
    "O banco de dados representa uma parte limitada\n"
    "do mundo real que interessa ao sistema.\n\n"
    "MUNDO REAL\n"
    "      ↓\n"
    "PARTE DA REALIDADE\n"
    "      ↓\n"
    "MINIMUNDO\n"
    "      ↓\n"
    "BANCO DE DADOS\n\n"
    "Exemplo: um sistema acadêmico pode trabalhar com\n"
    "alunos, cursos, disciplinas, matrículas e notas."
)

# 7 — BD, SGBD e SBD
add_slide(
    "BD × SGBD × SBD",
    "BANCO DE DADOS (BD)\n"
    "→ conjunto dos dados armazenados\n\n"
    "SGBD — Sistema de Gerenciamento de Banco de Dados\n"
    "→ software que permite definir, armazenar,\n"
    "   consultar, alterar e gerenciar os dados.\n\n"
    "SBD — Sistema de Banco de Dados\n"
    "→ Banco de Dados + SGBD"
)

# 8 — Analogia
add_slide(
    "Uma analogia: o armário",
    "ARMÁRIO → SGBD\n"
    "DOCUMENTOS → DADOS\n"
    "ORGANIZAÇÃO → ESTRUTURA\n\n"
    "Guardar documento → inserir\n"
    "Encontrar documento → consultar\n"
    "Modificar documento → atualizar\n"
    "Retirar documento → excluir\n\n"
    "O SGBD faz muito mais do que simplesmente guardar arquivos."
)

# 9 — CRUD
add_slide(
    "Operações básicas: CRUD",
    "C — Create → Criar / Inserir\n"
    "R — Read → Ler / Consultar\n"
    "U — Update → Atualizar / Alterar\n"
    "D — Delete → Excluir\n\n"
    "Exemplo:\n"
    "Cadastrar aluno → Consultar aluno → Alterar curso → Excluir cadastro"
)

# 10 — História
add_slide(
    "Um pouco de história",
    "• Final dos anos 1960: primeiros sistemas de bancos de dados\n"
    "• 1970: IBM publica trabalho sobre o modelo relacional\n"
    "• System R: projeto experimental da IBM\n"
    "• Desenvolvimento da SQL — Structured Query Language\n"
    "• SQL tornou-se padrão para bancos relacionais\n"
    "• System R evoluiu para tecnologias comerciais como SQL/DS e DB2\n\n"
    "Nesta aula: apenas uma visão geral. O foco é entender os conceitos."
)

# 11 — Tipos
add_slide(
    "Principais tipos de bancos de dados",
    "RELACIONAL\n"
    "→ dados organizados em tabelas\n"
    "→ linhas e colunas\n"
    "→ relacionamentos entre dados\n\n"
    "NoSQL\n"
    "→ diferentes modelos, como documentos,\n"
    "   chave-valor e grafos\n\n"
    "DISTRIBUÍDO\n"
    "→ dados gerenciados em múltiplos servidores/localizações"
)

# 12 — Atividade
add_activity_slide(
    "Atividade: imagine um sistema escolar",
    "1. Quais dados precisamos armazenar?\n\n"
    "2. Quais informações podemos ter sobre um aluno?\n\n"
    "3. O que o sistema precisa fazer com esses dados?\n"
    "   • Criar?\n"
    "   • Consultar?\n"
    "   • Atualizar?\n"
    "   • Excluir?\n\n"
    "4. Como vocês organizariam esses dados?"
)

# 13 — Fechamento
add_slide(
    "O que aprendemos hoje?",
    "• Dado e informação\n"
    "• Banco de Dados\n"
    "• Minimundo\n"
    "• SGBD\n"
    "• Sistema de Banco de Dados\n"
    "• CRUD\n"
    "• Visão geral dos tipos de BD\n\n"
    "IDEIA CENTRAL:\n"
    "Organizar dados para armazenar, recuperar e modificar informações."
)

# 14 — Gancho para a próxima aula
add_slide(
    "Próxima aula: como organizar esses dados?",
    "Hoje:\n"
    "REALIDADE → DADOS → BANCO DE DADOS\n\n"
    "Próxima etapa:\n"
    "BANCO DE DADOS RELACIONAL\n"
    "        ↓\n"
    "      TABELAS\n"
    "        ↓\n"
    "COLUNAS + LINHAS\n"
    "        ↓\n"
    "      CHAVES\n"
    "        ↓\n"
    "RELACIONAMENTOS\n\n"
    "E depois: SQL"
)

# Ajuste simples de tamanho das fontes
for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if shape == slide.shapes.title:
                    run.font.size = Pt(28)
                else:
                    run.font.size = Pt(20)

file_path = "/mnt/data/introducao_bancos_de_dados_aula_1.pptx"
prs.save(file_path)

file_path

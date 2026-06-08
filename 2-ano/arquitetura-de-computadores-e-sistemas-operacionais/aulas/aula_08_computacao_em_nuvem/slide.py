from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1 - Capa
add_slide(
    "Computação em Nuvem (IaaS, PaaS e SaaS)",
    "Conceitos, benefícios, aplicações e principais provedores de serviços em nuvem."
)

# Slide 2 - Introdução
add_slide(
    "O que é Computação em Nuvem?",
    "A computação em nuvem permite acessar recursos de TI pela internet.\n\n"
    "• Armazenamento de dados\n"
    "• Execução de aplicações\n"
    "• Compartilhamento de informações\n"
    "• Escalabilidade de serviços\n\n"
    "Está presente em empresas e no cotidiano das pessoas."
)

# Slide 3 - Benefícios da Computação em Nuvem
add_slide(
    "Principais Benefícios",
    "• Redução de custos com infraestrutura\n"
    "• Escalabilidade sob demanda\n"
    "• Alta disponibilidade dos serviços\n"
    "• Acesso remoto de qualquer lugar\n"
    "• Segurança e backup dos dados\n"
    "• Atualizações automáticas"
)

# Slide 4 - IaaS
add_slide(
    "Infraestrutura como Serviço (IaaS)",
    "Fornece recursos computacionais básicos como:\n\n"
    "• Servidores virtuais\n"
    "• Armazenamento\n"
    "• Redes\n\n"
    "A empresa gerencia o sistema operacional e aplicações, enquanto o provedor cuida da infraestrutura física."
)

# Slide 5 - Aplicações do IaaS
add_slide(
    "Exemplos de Uso do IaaS",
    "• Hospedagem de sites\n"
    "• Backup e recuperação de dados\n"
    "• Ambientes de teste e desenvolvimento\n"
    "• Expansão rápida de infraestrutura\n\n"
    "Permite criar ambientes completos sem comprar servidores físicos."
)

# Slide 6 - PaaS
add_slide(
    "Plataforma como Serviço (PaaS)",
    "Oferece uma plataforma pronta para desenvolvimento.\n\n"
    "Inclui:\n"
    "• Infraestrutura gerenciada\n"
    "• Banco de dados\n"
    "• Ferramentas de programação\n"
    "• Servidores de aplicação\n\n"
    "O foco do desenvolvedor fica no software."
)

# Slide 7 - Aplicações do PaaS
add_slide(
    "Exemplos de Uso do PaaS",
    "• Desenvolvimento de aplicativos móveis\n"
    "• Aplicações web\n"
    "• APIs e microsserviços\n"
    "• Projetos de startups\n\n"
    "Permite lançar produtos com mais rapidez."
)

# Slide 8 - SaaS
add_slide(
    "Software como Serviço (SaaS)",
    "O software é disponibilizado pela internet.\n\n"
    "Características:\n"
    "• Não exige instalação local\n"
    "• Atualizações automáticas\n"
    "• Pagamento por assinatura\n\n"
    "Exemplos: Google Workspace, Microsoft 365 e Salesforce."
)

# Slide 9 - Principais Players do Mercado
add_slide(
    "Big Three da Computação em Nuvem",
    "Amazon Web Services (AWS)\n"
    "• Grande variedade de serviços\n\n"
    "Microsoft Azure\n"
    "• Forte integração com produtos Microsoft\n\n"
    "Google Cloud Platform (GCP)\n"
    "• Destaque em Big Data, IA e Machine Learning"
)

# Slide 10 - Comparação dos Provedores
add_slide(
    "AWS x Azure x GCP",
    "AWS:\n"
    "• Ampla gama de serviços\n"
    "• Clientes: Netflix, Airbnb\n\n"
    "Azure:\n"
    "• Integração corporativa\n"
    "• Clientes: LinkedIn, Adobe\n\n"
    "GCP:\n"
    "• Big Data e Inteligência Artificial\n"
    "• Clientes: Spotify, HSBC"
)

# Slide 11 - Conclusão
add_slide(
    "Conclusão",
    "A computação em nuvem revolucionou a área de TI.\n\n"
    "• IaaS oferece infraestrutura.\n"
    "• PaaS oferece plataforma para desenvolvimento.\n"
    "• SaaS entrega software pronto para uso.\n\n"
    "Esses modelos tornam a tecnologia mais acessível, escalável e econômica."
)

# Slide 12 - Atividade
add_activity_slide(
    "Atividade",
    "1. Qual a principal diferença entre IaaS, PaaS e SaaS?\n\n"
    "2. Cite dois benefícios da computação em nuvem.\n\n"
    "3. Qual provedor é mais conhecido por Big Data e IA?\n\n"
    "4. Dê um exemplo de aplicação para IaaS."
)

file_path = "Computacao_em_Nuvem_IaaS_PaaS_SaaS.pptx"
prs.save(file_path)

print(file_path)
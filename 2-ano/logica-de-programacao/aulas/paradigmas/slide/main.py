from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1 - Capa
add_slide(
    "Paradigmas de Programação",
    "Conceitos, características e principais tipos de paradigmas utilizados no desenvolvimento de software."
)

# Slide 2 - Definição
add_slide(
    "Definição",
    "Um paradigma de programação é um modelo que orienta a forma como programas são desenvolvidos e estruturados."
)

# Slide 3 - Importância
add_slide(
    "Importância dos Paradigmas",
    "Fornecem padrões para organizar códigos, facilitando manutenção, reutilização e compreensão dos sistemas."
)

# Slide 4 - Linguagens Multiparadigma
add_slide(
    "Linguagens Multiparadigma",
    "Python é considerada uma linguagem multiparadigma por permitir programação estruturada, funcional e orientada a objetos."
)

# Slide 5 - Principais Paradigmas
add_slide(
    "Tipos de Paradigmas",
    "Estruturada, Procedural, Interativa, Funcional e Orientada a Objetos."
)

# Slide 6 - Programação Estruturada
add_slide(
    "Programação Estruturada",
    "Baseada na execução sequencial de instruções usando estruturas condicionais, laços de repetição e funções."
)

# Slide 7 - Programação Procedural
add_slide(
    "Programação Procedural",
    "Organiza o programa em procedimentos e funções para facilitar a reutilização e modularização do código."
)

# Slide 8 - Programação Interativa
add_slide(
    "Programação Interativa",
    "Permite executar comandos e obter respostas imediatas. Exemplos: Python REPL e Jupyter Notebook."
)

# Slide 9 - Programação Funcional
add_slide(
    "Programação Funcional",
    "Baseada no uso de funções, recursão, expressões lambda, map(), filter() e list comprehensions."
)

# Slide 10 - Programação Orientada a Objetos
add_slide(
    "Programação Orientada a Objetos",
    "Organiza programas em classes e objetos contendo atributos e métodos."
)

# Slide 11 - Vantagens da POO
add_slide(
    "Vantagens da POO",
    "Promove encapsulamento, herança, polimorfismo, reutilização de código e escalabilidade."
)

# Slide 12 - Conclusão
add_slide(
    "Conclusão",
    "Cada paradigma possui características específicas. Python se destaca por oferecer suporte a múltiplos paradigmas."
)

# Slide 13 - Atividade
add_activity_slide(
    "Atividade",
    "1. O que é um paradigma de programação?\n"
    "2. Qual a diferença entre programação estruturada e funcional?\n"
    "3. Cite duas vantagens da programação orientada a objetos.\n"
    "4. Por que Python é considerada multiparadigma?"
)

file_path = "paradigmas_programacao.pptx"
prs.save(file_path)

file_path
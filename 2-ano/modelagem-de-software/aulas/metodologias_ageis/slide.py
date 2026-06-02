from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Metodologias Ágeis e Engenharia de Software"
slide.placeholders[1].text = "Scrum, Kanban e Desenvolvimento Ágil"

# Slide 2
add_slide(
    "Objetivos da Aula",
    "• Entender os princípios do Manifesto Ágil\n"
    "• Conhecer Scrum e Kanban\n"
    "• Aplicar conceitos em atividades práticas\n"
    "• Desenvolver colaboração e organização em projetos"
)

# Slide 3
add_slide(
    "O que são Metodologias Ágeis?",
    "Abordagens que ajudam equipes a desenvolver software de forma "
    "mais rápida, colaborativa e flexível.\n\n"
    "Foco em entregas frequentes e adaptação às mudanças."
)

# Slide 4
add_slide(
    "Manifesto Ágil",
    "Valores principais:\n"
    "• Pessoas e interações\n"
    "• Software funcionando\n"
    "• Colaboração com o cliente\n"
    "• Resposta rápida às mudanças"
)

# Slide 5
add_slide(
    "Por que usar Metodologias Ágeis?",
    "• Maior flexibilidade\n"
    "• Melhor comunicação da equipe\n"
    "• Entregas frequentes\n"
    "• Menos retrabalho\n"
    "• Maior satisfação do cliente"
)

# Slide 6
add_slide(
    "Scrum",
    "Metodologia baseada em ciclos curtos chamados Sprints.\n\n"
    "Cada Sprint produz uma parte funcional do sistema."
)

# Slide 7
add_slide(
    "Papéis do Scrum",
    "Product Owner:\nDefine prioridades.\n\n"
    "Scrum Master:\nRemove impedimentos e facilita o processo.\n\n"
    "Equipe de Desenvolvimento:\nConstrói o produto."
)

# Slide 8
add_slide(
    "Eventos do Scrum",
    "• Sprint Planning\n"
    "• Daily Scrum\n"
    "• Sprint Review\n"
    "• Sprint Retrospective"
)

# Slide 9
add_slide(
    "Artefatos do Scrum",
    "• Product Backlog\n"
    "• Sprint Backlog\n"
    "• Incremento do Produto"
)

# Slide 10
add_slide(
    "Kanban",
    "Método visual para organizar tarefas e acompanhar o fluxo de trabalho."
)

# Slide 11
add_slide(
    "Princípios do Kanban",
    "• Visualização do trabalho\n"
    "• Limite de trabalho em progresso (WIP)\n"
    "• Fluxo contínuo\n"
    "• Melhoria contínua"
)

# Slide 12
add_slide(
    "Exemplo de Quadro Kanban",
    "A Fazer\n"
    "↓\n"
    "Em Progresso\n"
    "↓\n"
    "Concluído"
)

# Slide 13
add_slide(
    "Scrum x Kanban",
    "SCRUM:\n"
    "• Trabalha com Sprints\n"
    "• Papéis definidos\n\n"
    "KANBAN:\n"
    "• Fluxo contínuo\n"
    "• Mais flexível"
)

# Slide 14
add_slide(
    "Desenvolvimento Iterativo e Incremental",
    "Exemplo:\n"
    "1. Cadastro de Pacientes\n"
    "2. Agendamento de Consultas\n"
    "3. Cancelamento de Consultas\n\n"
    "Cada etapa adiciona novas funcionalidades."
)

# Slide 15
add_slide(
    "Entrega Contínua",
    "Após cada melhoria:\n"
    "• Testar\n"
    "• Validar\n"
    "• Entregar ao usuário\n\n"
    "Resultado: mais qualidade e rapidez."
)

# Slide 16
add_slide(
    "Planejamento Ágil",
    "• Divisão do trabalho em Sprints\n"
    "• Priorização das tarefas\n"
    "• Uso de Story Points para estimar esforço"
)

# Slide 17
add_slide(
    "Ferramentas Utilizadas",
    "• Trello\n"
    "• Jira\n"
    "• Miro\n"
    "• GitLab\n"
    "• Excel"
)

# Slide 18
add_slide(
    "Atividade Prática 1",
    "Simulação de Sprint Planning\n\n"
    "Criar um backlog para um sistema de biblioteca.\n"
    "Definir prioridades e tarefas do Sprint."
)

# Slide 19
add_slide(
    "Atividade Prática 2",
    "Criar um quadro Kanban.\n\n"
    "Organizar tarefas nas colunas:\n"
    "• A Fazer\n"
    "• Em Progresso\n"
    "• Concluído"
)

# Slide 20
add_slide(
    "Perguntas para Reflexão",
    "1. Como Scrum e Kanban podem melhorar a produtividade?\n\n"
    "2. Quando o Kanban pode ser mais vantajoso que o Scrum?"
)

# Slide 21
add_slide(
    "Exemplo do Dia a Dia",
    "Planejamento de uma viagem:\n"
    "• Definir tarefas\n"
    "• Organizar prioridades\n"
    "• Acompanhar progresso\n\n"
    "Assim funcionam as metodologias ágeis."
)

# Slide 22
add_slide(
    "Conclusão",
    "Metodologias Ágeis ajudam equipes a:\n"
    "• Trabalhar melhor em grupo\n"
    "• Entregar valor rapidamente\n"
    "• Adaptar-se às mudanças\n"
    "• Melhorar continuamente"
)

prs.save("Metodologias_Ageis.pptx")

print("Apresentação criada com sucesso!")
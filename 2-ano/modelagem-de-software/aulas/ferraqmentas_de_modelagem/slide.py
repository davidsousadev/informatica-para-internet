from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1 - Capa

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Ferramentas e Técnicas de Modelagem"
slide.placeholders[1].text = "Aula de Desenvolvimento de Sistemas\nEnsino Médio"

# Slide 2

add_slide(
"Competências da Aula",
"• Conhecer linguagens de modelagem como UML.\n"
"• Identificar ferramentas de modelagem de software.\n"
"• Criar diagramas UML e protótipos.\n"
"• Escolher a ferramenta adequada para cada projeto.\n"
"• Melhorar a comunicação e o planejamento de sistemas."
)

# Slide 3

add_slide(
"Objetivos da Aula",
"• Apresentar Lucidchart, Enterprise Architect, Figma e Adobe XD.\n"
"• Demonstrar a criação de diagramas UML.\n"
"• Desenvolver protótipos interativos.\n"
"• Comparar vantagens e limitações das ferramentas."
)

# Slide 4

add_slide(
"O que é Modelagem?",
"A modelagem é a representação visual de um sistema antes da sua construção.\n\n"
"Ela ajuda a planejar, organizar ideias e facilitar a comunicação entre equipes."
)

# Slide 5

add_slide(
"Por que Utilizar Ferramentas de Modelagem?",
"• Organizam requisitos.\n"
"• Facilitam a documentação.\n"
"• Melhoram a comunicação.\n"
"• Reduzem erros de desenvolvimento.\n"
"• Aumentam a produtividade."
)

# Slide 6

add_slide(
"Lucidchart",
"Ferramenta online para criação de diagramas UML.\n\n"
"Características:\n"
"• Fácil de usar.\n"
"• Colaboração em tempo real.\n"
"• Diagramas de casos de uso, classes e sequência.\n\n"
"Ideal para equipes e projetos acadêmicos."
)

# Slide 7

add_slide(
"Enterprise Architect",
"Ferramenta profissional para modelagem UML.\n\n"
"Características:\n"
"• Modelagem avançada.\n"
"• Controle de versões.\n"
"• Documentação detalhada.\n\n"
"Ideal para projetos grandes e complexos."
)

# Slide 8

add_slide(
"Figma",
"Ferramenta colaborativa de design e prototipação.\n\n"
"Características:\n"
"• Criação de interfaces.\n"
"• Trabalho em equipe online.\n"
"• Protótipos interativos.\n\n"
"Muito utilizada por designers UX/UI."
)

# Slide 9

add_slide(
"Adobe XD",
"Software para design de interfaces e protótipos.\n\n"
"Características:\n"
"• Simulação de navegação.\n"
"• Protótipos de alta fidelidade.\n"
"• Integração com ferramentas Adobe."
)

# Slide 10

add_slide(
"Comparação das Ferramentas",
"Lucidchart → Diagramas simples e colaboração.\n\n"
"Enterprise Architect → Projetos complexos.\n\n"
"Figma → Design de interfaces.\n\n"
"Adobe XD → Prototipagem avançada."
)

# Slide 11

add_slide(
"Exemplo de Diagrama de Classes",
"Sistema de Reserva de Hotel:\n\n"
"Cliente\n"
"Reserva\n"
"Quarto\n"
"Pagamento\n\n"
"Essas classes se relacionam para representar o funcionamento do sistema."
)

# Slide 12

add_slide(
"Exemplo de Caso de Uso",
"Sistema de Biblioteca:\n\n"
"Ator: Aluno\n\n"
"Casos de Uso:\n"
"• Consultar livros\n"
"• Realizar empréstimo\n"
"• Devolver livro\n"
"• Renovar empréstimo"
)

# Slide 13

add_slide(
"Exemplo de Protótipo",
"Aplicativo de Pedidos de Comida:\n\n"
"• Tela de Login\n"
"• Tela de Menu\n"
"• Carrinho\n"
"• Finalização do Pedido\n\n"
"Ferramenta sugerida: Figma."
)

# Slide 14

add_slide(
"Atividade Prática 1",
"Criar um Diagrama de Classes para um Sistema de Reservas de Hotel.\n\n"
"Classes sugeridas:\n"
"Cliente, Reserva, Quarto e Pagamento."
)

# Slide 15

add_slide(
"Atividade Prática 2",
"Criar um Protótipo Interativo no Figma.\n\n"
"Desenvolver as telas:\n"
"• Login\n"
"• Cardápio\n"
"• Carrinho\n"
"• Finalização do Pedido"
)

# Slide 16

add_slide(
"Atividade Prática 3",
"Criar um Diagrama de Casos de Uso para um Sistema de Biblioteca.\n\n"
"Identificar atores e funcionalidades."
)

# Slide 17

add_slide(
"Perguntas para Reflexão",
"1. Qual ferramenta é melhor para trabalho colaborativo?\n\n"
"2. Como protótipos ajudam no desenvolvimento?\n\n"
"3. Em quais situações usar UML?"
)

# Slide 18

add_slide(
"Mercado de Trabalho",
"Profissionais de Desenvolvimento de Sistemas utilizam diariamente ferramentas de modelagem.\n\n"
"Conhecer UML, Figma e ferramentas colaborativas é um diferencial para estágios e empregos."
)

# Slide 19

add_slide(
"Resumo da Aula",
"• Ferramentas de modelagem facilitam o desenvolvimento.\n"
"• UML ajuda a documentar sistemas.\n"
"• Figma e Adobe XD permitem criar protótipos.\n"
"• A escolha da ferramenta depende do projeto."
)

# Slide 20 - Encerramento

add_slide(
"Obrigado!",
"Dúvidas?\n\nVamos praticar a modelagem de software!"
)

arquivo = "Ferramentas_Tecnicas_Modelagem.pptx"
prs.save(arquivo)

print(f'Apresentação salva como: {arquivo}')

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Slide inicial
add_title_slide("A Revolução das Telas", "Introdução ao UI/UX Design")

# Slides
add_slide("O que é a Internet", [
    "Rede mundial de computadores",
    "Conecta dispositivos desde os anos 1990",
    "Facilita comunicação e compartilhamento"
])

add_slide("A Revolução das Telas", [
    "Mudança no acesso à internet",
    "Celulares, desktops, laptops, tablets e TVs"
])

add_slide("Design Responsivo", [
    "Sites se adaptam ao tamanho da tela",
    "Importante para múltiplos dispositivos"
])

add_slide("Pixel", [
    "Menor unidade de imagem",
    "Define nitidez e qualidade"
])

add_slide("Resoluções", [
    "HD: 1280 x 720",
    "Full HD: 1920 x 1080",
    "2K e 4K"
])

add_slide("Densidade de Pixels", [
    "Medida em ppi",
    "Mais ppi = mais qualidade"
])

add_slide("Interfaces", [
    "Conectam usuário ao sistema",
    "Simplificam o uso"
])

add_slide("Consistência", [
    "Interface deve agir como esperado",
    "Falhas reduzem confiança"
])

add_slide("Feedback", [
    "Resposta do sistema ao usuário",
    "Ex: mensagens e sons"
])

# Salvar
file_path = "revolucao_das_telas.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
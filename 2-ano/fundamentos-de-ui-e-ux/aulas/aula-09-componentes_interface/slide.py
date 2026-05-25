from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1
add_slide(
    "Aula 09 - Componentes de Interface",
    "Fundamentos de UI e UX\nEstruturas e componentes utilizados em aplicações web."
)

# Slide 2
add_slide(
    "Objetivos da Aula",
    "• Conhecer a estrutura básica de um site\n"
    "• Identificar componentes comuns de interface\n"
    "• Entender quando utilizar cada componente\n"
    "• Aplicar boas práticas de usabilidade"
)

# Slide 3
add_slide(
    "Estrutura Básica de um Site",
    "Uma interface web normalmente é composta por:\n"
    "• Cabeçalho (Header)\n"
    "• Seção Principal (Main)\n"
    "• Rodapé (Footer)"
)

# Slide 4
add_slide(
    "Cabeçalho (Header)",
    "• Logo ou identidade visual\n"
    "• Menu de navegação\n"
    "• Barra de busca\n"
    "• Login e serviços\n"
    "• Principal ponto de atenção do usuário"
)

# Slide 5
add_slide(
    "Cabeçalho Responsivo",
    "• Adaptação para dispositivos móveis\n"
    "• Uso do menu sanduíche (hamburger)\n"
    "• Busca simplificada\n"
    "• Melhor aproveitamento do espaço"
)

# Slide 6
add_slide(
    "Rodapé (Footer)",
    "• Informações institucionais\n"
    "• Política de privacidade\n"
    "• Redes sociais\n"
    "• Contato\n"
    "• SEO e informações complementares"
)

# Slide 7
add_slide(
    "Seção Principal (Main)",
    "• Conteúdo principal da página\n"
    "• Área dinâmica\n"
    "• Informações relevantes ao usuário\n"
    "• Pode conter menus laterais e destaques"
)

# Slide 8
add_slide(
    "Elementos Comuns",
    "• Textos\n"
    "• Links e Botões\n"
    "• Ícones\n"
    "• Imagens e Vetores\n"
    "• Formulários\n"
    "• Tabelas"
)

# Slide 9
add_slide(
    "Links e Botões",
    "• Devem ser facilmente identificáveis\n"
    "• Estados: normal, hover, pressionado e desabilitado\n"
    "• Facilitam a navegação e interação"
)

# Slide 10
add_slide(
    "Ícones, Imagens e Vetores",
    "• Complementam informações visuais\n"
    "• Melhoram a compreensão\n"
    "• Vetores (SVG) possuem alta qualidade\n"
    "• Devem possuir acessibilidade"
)

# Slide 11
add_slide(
    "Formulários e Tabelas",
    "• Captura de dados do usuário\n"
    "• Campos de texto, seleção e upload\n"
    "• Validação visual\n"
    "• Tabelas para dados tabulares"
)

# Slide 12
add_slide(
    "Alertas e Menus",
    "• Alertas de sucesso, erro e aviso\n"
    "• Menu Sanfona (Accordion)\n"
    "• Efeito Colapso\n"
    "• Menu Dropdown"
)

# Slide 13
add_slide(
    "Abas e Paginação",
    "• Organização de conteúdos\n"
    "• Navegação entre seções\n"
    "• Redução da carga visual\n"
    "• Melhora da experiência do usuário"
)

# Slide 14
add_slide(
    "Carrossel e Modal",
    "• Carrossel para múltiplas imagens\n"
    "• Modal para informações em destaque\n"
    "• Evita mudança de página\n"
    "• Foco na interação"
)

# Slide 15
add_slide(
    "Offcanvas, Toasts e Tooltips",
    "• Offcanvas: barra lateral\n"
    "• Toast: notificações rápidas\n"
    "• Tooltip: dica ao passar o mouse\n"
    "• Popover: informações adicionais ao clicar"
)

# Slide 16
add_slide(
    "Barra de Progresso e Cards",
    "• Acompanhamento de processos\n"
    "• Exibição visual de progresso\n"
    "• Cards organizam conteúdos de forma modular"
)

# Slide 17
add_slide(
    "Breadcrumbs (Trilhas de Navegação)",
    "• Indicam a localização do usuário\n"
    "• Facilitam a navegação\n"
    "• Muito utilizados em sites grandes"
)

# Slide 18
add_slide(
    "Bootstrap e Componentes Prontos",
    "• Framework front-end popular\n"
    "• Componentes prontos e responsivos\n"
    "• Facilita o desenvolvimento\n"
    "• Documentação oficial: getbootstrap.com"
)

# Slide 19
add_slide(
    "Boas Práticas de UI",
    "• Consistência visual\n"
    "• Clareza na navegação\n"
    "• Feedback ao usuário\n"
    "• Responsividade\n"
    "• Acessibilidade"
)

# Slide 20
add_activity_slide(
    "Atividade",
    "Imagine uma aplicação para sua área.\n\n"
    "1. Quais componentes seriam necessários?\n"
    "2. Onde ficariam header, main e footer?\n"
    "3. Quais componentes melhorariam a experiência do usuário?"
)

file_path = "aula_09_componentes_de_interface.pptx"
prs.save(file_path)

print(file_path)
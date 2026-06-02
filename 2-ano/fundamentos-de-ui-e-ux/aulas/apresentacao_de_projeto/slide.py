from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1 - Capa
add_slide(
    "Apresentação do Projeto",
    "Design de Interfaces, Stakeholders e Construção de Portfólio"
)

# Slide 2 - Processo de Design
add_slide(
    "Processo de Design de Interfaces",
    "Segundo Preece, Rogers e Sharp (2002), o processo envolve:\n"
    "• Análise de requisitos\n"
    "• Desenvolvimento de designs alternativos\n"
    "• Prototipação\n"
    "• Avaliação\n\n"
    "O processo é cíclico e pode retornar às etapas anteriores para melhorias."
)

# Slide 3 - Evolução do Projeto
add_slide(
    "Evolução Contínua do Projeto",
    "Mesmo após a conclusão e lançamento, uma aplicação web deve continuar evoluindo.\n"
    "Quando os objetivos iniciais são atingidos, o projeto segue para desenvolvimento e lançamento."
)

# Slide 4 - Quem são os Stakeholders?
add_slide(
    "Stakeholders",
    "São pessoas ou entidades interessadas ou afetadas pelo projeto:\n"
    "• Clientes\n"
    "• Gerentes\n"
    "• Proprietários\n"
    "• Fornecedores\n"
    "• Empresas parceiras\n"
    "• Equipes de desenvolvimento"
)

# Slide 5 - Apresentação aos Stakeholders
add_slide(
    "Apresentação do Projeto",
    "O designer deve apresentar:\n"
    "• Documentação do design\n"
    "• Protótipos desenvolvidos\n"
    "• Justificativas para as escolhas da interface\n"
    "• Benefícios das soluções adotadas"
)

# Slide 6 - Documentação do Projeto
add_slide(
    "Documentação do Design",
    "Registro formal do projeto contendo:\n"
    "• Guia de identidade visual\n"
    "• Paleta de cores\n"
    "• Tipografias\n"
    "• Logos\n"
    "• Elementos gráficos\n"
    "• Estrutura e seções da interface"
)

# Slide 7 - Justificando Decisões
add_slide(
    "Justificando as Escolhas",
    "Durante a apresentação, é importante:\n"
    "• Explicar as decisões tomadas\n"
    "• Demonstrar como elas beneficiam o usuário\n"
    "• Defender estratégias importantes para o projeto\n"
    "• Avaliar sugestões sem comprometer a usabilidade"
)

# Slide 8 - Uso de Protótipos
add_slide(
    "Protótipos no Processo",
    "• Apresentar protótipos de baixa fidelidade primeiro\n"
    "• Evoluir para protótipos de alta fidelidade\n"
    "• Realizar reuniões frequentes com stakeholders\n"
    "• Garantir alinhamento durante todo o projeto"
)

# Slide 9 - Comunicação Eficaz
add_slide(
    "Boas Práticas de Comunicação",
    "• Utilizar linguagem clara e simples\n"
    "• Evitar jargões técnicos\n"
    "• Escutar feedbacks\n"
    "• Estar aberto a sugestões\n"
    "• Saber dizer 'não' quando necessário para preservar a qualidade do design"
)

# Slide 10 - Construindo um Portfólio
add_slide(
    "Construindo seu Portfólio",
    "Um portfólio é uma coleção organizada dos trabalhos realizados.\n"
    "Ele demonstra habilidades, metodologias e resultados alcançados em projetos de design."
)

# Slide 11 - Conteúdo do Portfólio
add_slide(
    "O que Incluir no Portfólio?",
    "• Projetos desenvolvidos\n"
    "• Processo criativo\n"
    "• Ferramentas utilizadas\n"
    "• Wireframes e mockups\n"
    "• Testes de usabilidade\n"
    "• Aplicação dos princípios de UX e UI"
)

# Slide 12 - Conclusão
add_slide(
    "Conclusão",
    "A apresentação aos stakeholders garante alinhamento e validação do projeto.\n"
    "Já o portfólio evidencia a experiência e o impacto do profissional na criação de soluções centradas no usuário."
)

file_path = "Apresentacao_Projeto_Stakeholders_Portfolio.pptx"
prs.save(file_path)

print(file_path)
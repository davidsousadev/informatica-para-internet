from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1
add_slide(
    "Prototipagem e Design de Interfaces",
    "Arquitetura da Informação, Wireframes, Mockups, Protótipos e uso do Figma."
)

# Slide 2
add_slide(
    "Objetivos da Prototipagem",
    "• Transformar requisitos em soluções visuais\n"
    "• Testar ideias antes do desenvolvimento\n"
    "• Reduzir erros e retrabalho\n"
    "• Melhorar a experiência do usuário"
)

# Slide 3
add_slide(
    "Arquitetura da Informação",
    "Organiza o conteúdo de forma clara e intuitiva.\n\n"
    "Exemplo de estrutura:\n"
    "Página Inicial\n"
    "├─ Artistas\n"
    "├─ Álbuns\n"
    "├─ Planos\n"
    "├─ Assinar\n"
    "└─ Login"
)

# Slide 4
add_slide(
    "Etapas da Prototipagem",
    "1. Wireframe (esboço)\n"
    "2. Mockup (design visual)\n"
    "3. Protótipo (interativo)\n"
    "4. Ferramenta final"
)

# Slide 5
add_slide(
    "Wireframe",
    "Protótipo de baixa fidelidade.\n"
    "Utiliza formas simples para definir layout e posicionamento dos elementos.\n"
    "Pode ser feito em papel ou ferramentas digitais."
)

# Slide 6
add_slide(
    "Mockup",
    "Protótipo de alta fidelidade.\n"
    "Inclui cores, tipografia, imagens e identidade visual.\n"
    "Representa a aparência final da interface."
)

# Slide 7
add_slide(
    "Protótipo Interativo",
    "Permite navegação entre telas.\n"
    "Simula o comportamento do sistema antes da implementação.\n"
    "Facilita testes de usabilidade."
)

# Slide 8
add_slide(
    "Ferramentas de Prototipagem",
    "• Figma\n"
    "• Adobe XD\n\n"
    "Ambas permitem criar wireframes, mockups e protótipos interativos."
)

# Slide 9
add_slide(
    "Estudo de Caso: Editora UI/UX",
    "Empresa fictícia de livros acadêmicos.\n"
    "Objetivo: desenvolver a interface de um site institucional."
)

# Slide 10
add_slide(
    "Wireframe do Projeto",
    "Estrutura da página:\n"
    "• Cabeçalho\n"
    "• Seção principal (Hero)\n"
    "• Cards de novidades\n"
    "• Rodapé"
)

# Slide 11
add_slide(
    "Construção no Figma",
    "• Criar Frame\n"
    "• Configurar Grid de 12 colunas\n"
    "• Adicionar componentes\n"
    "• Organizar layout responsivo"
)

# Slide 12
add_slide(
    "Criando Protótipos Interativos",
    "• Duplicar frames\n"
    "• Criar páginas secundárias\n"
    "• Configurar interações\n"
    "• Navegação entre telas"
)

# Slide 13
add_slide(
    "Versão Mobile",
    "Adaptação para smartphones:\n"
    "• Layout vertical\n"
    "• Cards empilhados\n"
    "• Menu sanduíche\n"
    "• Interface responsiva"
)

# Slide 14
add_slide(
    "Menu Sanduíche",
    "Componente utilizado em telas pequenas.\n"
    "Permite abrir e fechar menus através de sobreposição (overlay)."
)

# Slide 15
add_slide(
    "Conclusão",
    "• Arquitetura da Informação organiza conteúdos.\n"
    "• Wireframes definem estrutura.\n"
    "• Mockups representam o visual.\n"
    "• Protótipos validam a experiência antes do desenvolvimento."
)

file_path = "Prototipagem_Figma.pptx"
prs.save(file_path)

print(file_path)
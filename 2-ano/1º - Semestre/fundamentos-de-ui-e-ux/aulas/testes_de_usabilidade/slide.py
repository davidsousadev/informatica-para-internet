from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1
add_slide(
    "Testes de Usabilidade",
    "Avaliação da facilidade de uso de interfaces para identificar problemas e melhorar a experiência do usuário."
)

# Slide 2
add_slide(
    "Por que Testar a Usabilidade?",
    "• Desenvolvedores conhecem demais o sistema\n"
    "• Problemas podem passar despercebidos\n"
    "• Steve Krug: 'Se você quiser um ótimo site, deve testá-lo'\n"
    "• Melhora a experiência do usuário"
)

# Slide 3
add_slide(
    "Formas de Avaliação",
    "1. Avaliação com Especialistas\n"
    "2. Testes com Usuários\n\n"
    "Ambas são complementares e importantes no desenvolvimento de interfaces."
)

# Slide 4
add_slide(
    "Avaliação Heurística",
    "Método de inspeção realizado por especialistas.\n"
    "Baseia-se nas 10 Heurísticas de Nielsen (1994).\n"
    "Permite identificar problemas de usabilidade com baixo custo."
)

# Slide 5
add_slide(
    "As 10 Heurísticas de Nielsen",
    "1. Visibilidade do status do sistema\n"
    "2. Compatibilidade com o mundo real\n"
    "3. Controle e liberdade do usuário\n"
    "4. Consistência e padrões\n"
    "5. Prevenção de erros"
)

# Slide 6
add_slide(
    "As 10 Heurísticas de Nielsen (Continuação)",
    "6. Reconhecimento em vez de lembrança\n"
    "7. Flexibilidade e eficiência de uso\n"
    "8. Estética e design minimalista\n"
    "9. Diagnóstico e correção de erros\n"
    "10. Ajuda e documentação"
)

# Slide 7
add_slide(
    "Como Aplicar as Heurísticas",
    "• Compreender o produto\n"
    "• Avaliar cada heurística\n"
    "• Identificar problemas\n"
    "• Propor soluções\n"
    "• Discutir resultados com a equipe"
)

# Slide 8
add_slide(
    "Exemplo: Visibilidade do Status",
    "Problema: usuário adiciona produto ao carrinho e não recebe feedback.\n\n"
    "Solução: exibir mensagem de confirmação e atualizar contador do carrinho."
)

# Slide 9
add_slide(
    "Exemplo: Prevenção de Erros",
    "Problema: cadastro permite envio sem CPF obrigatório.\n\n"
    "Solução: validar o campo antes do envio e exibir mensagens de erro."
)

# Slide 10
add_slide(
    "Testes com Usuários",
    "• Complementam as avaliações heurísticas\n"
    "• Observam o comportamento real dos usuários\n"
    "• Permitem descobrir problemas não previstos pelos especialistas"
)

# Slide 11
add_slide(
    "Características dos Testes",
    "• Ambiente controlado\n"
    "• Facilitador acompanha o processo\n"
    "• Técnica Thinking Aloud\n"
    "• Gravação de tela, áudio e vídeo"
)

# Slide 12
add_slide(
    "Tipos de Testes de Usabilidade",
    "• Testes A/B\n"
    "• Testes moderados e não moderados\n"
    "• Testes com protótipos\n"
    "• Testes de cenários\n"
    "• Eye Tracking"
)

# Slide 13
add_slide(
    "Teste A/B",
    "Comparação entre duas versões de uma interface.\n"
    "Avalia desempenho, preferência, tempo de execução e número de cliques."
)

# Slide 14
add_slide(
    "Teste de Cenários",
    "Usuários realizam tarefas em situações simuladas.\n"
    "Objetivo: identificar dificuldades e oportunidades de melhoria."
)

# Slide 15
add_slide(
    "Eye Tracking",
    "Rastreamento ocular para identificar áreas de maior atenção.\n"
    "Resultados apresentados em mapas de calor e trajetórias visuais."
)

# Slide 16
add_slide(
    "Teste de Usabilidade de Krug",
    "Método simples e econômico.\n"
    "• 3 usuários\n"
    "• Menos de 1 hora\n"
    "• Realizado mensalmente\n"
    "• Gravação da interação"
)

# Slide 17
add_slide(
    "Etapas do Teste de Krug",
    "1. Boas-vindas\n"
    "2. Perguntas pessoais\n"
    "3. Apresentação do site\n"
    "4. Execução de tarefas\n"
    "5. Sondagem\n"
    "6. Encerramento"
)

# Slide 18
add_slide(
    "Análise dos Resultados",
    "• Equipe assiste às gravações\n"
    "• Cada participante identifica problemas\n"
    "• Seleção dos 10 problemas mais críticos\n"
    "• Priorização das correções"
)

# Slide 19
add_slide(
    "Benefícios dos Testes de Usabilidade",
    "• Melhor experiência do usuário\n"
    "• Redução de erros\n"
    "• Aumento da satisfação\n"
    "• Economia de custos futuros\n"
    "• Maior qualidade do produto"
)

# Slide 20
add_slide(
    "Conclusão",
    "Testes de usabilidade são essenciais para criar interfaces eficientes, intuitivas e centradas no usuário. Avaliações heurísticas e testes com usuários devem ser usados em conjunto."
)

file_path = "Testes_de_Usabilidade.pptx"
prs.save(file_path)

print(file_path)
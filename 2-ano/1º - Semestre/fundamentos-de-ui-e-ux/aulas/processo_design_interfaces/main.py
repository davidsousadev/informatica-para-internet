from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Slide inicial
add_title_slide(
    "Processo de Design de Interfaces de Usuário",
    "Introdução ao UI/UX"
)

# Slides
add_slide("Objetivos da Aula", [
    "Compreender o processo de design de interfaces",
    "Entender a análise de requisitos",
    "Conhecer pesquisa com usuários",
    "Aplicar personas e jornada"
])

add_slide("Processo de Design", [
    "Dividido em 4 etapas principais",
    "Análise de requisitos",
    "Designs alternativos",
    "Prototipação e avaliação",
    "Processo cíclico"
])

add_slide("Análise de Requisitos", [
    "Identificar público-alvo",
    "Entender necessidades dos usuários",
    "Definir objetivos do sistema",
    "Base para todo o projeto"
])

add_slide("Designs Alternativos", [
    "Criação de diferentes soluções",
    "Design conceitual (funções)",
    "Design físico (interface visual)",
    "Cores, menus e layout"
])

add_slide("Prototipação", [
    "Versão inicial do sistema",
    "Permite interação do usuário",
    "Não precisa estar completa",
    "Usada para testes"
])

add_slide("Avaliação", [
    "Mede usabilidade do sistema",
    "Número de erros",
    "Facilidade de uso",
    "Satisfação do usuário"
])

add_slide("Processo Cíclico", [
    "Não é linear",
    "Pode voltar etapas anteriores",
    "Avaliação pode gerar novos requisitos"
])

add_slide("Pesquisa com Usuários", [
    "Entender necessidades reais",
    "Pode ocorrer em todo o projeto",
    "Base para decisões de design"
])

add_slide("Técnicas de Coleta", [
    "Entrevistas",
    "Questionários",
    "Grupos de foco",
    "Estudos de campo"
])

add_slide("Entrevistas", [
    "Conversas com usuários",
    "Podem ser estruturadas ou não",
    "Semiestruturadas são mais comuns",
    "Geram dados qualitativos"
])

add_slide("Questionários", [
    "Formulários com perguntas",
    "Podem ser abertas ou fechadas",
    "Geram dados quantitativos",
    "Fáceis de analisar"
])

add_slide("Personas", [
    "Usuários fictícios baseados em dados reais",
    "Representam o público-alvo",
    "Ex: idade, profissão, objetivos"
])

add_slide("Jornada do Usuário", [
    "Caminho do usuário no sistema",
    "Mostra interações e dificuldades",
    "Ajuda a identificar problemas"
])

add_slide("Importância", [
    "Melhorar experiência do usuário",
    "Criar interfaces eficientes",
    "Reduzir erros e frustrações"
])

add_slide("Atividade em Sala", [
    "Criar um aplicativo fictício",
    "Definir objetivo",
    "Criar persona",
    "Descrever jornada do usuário"
])

# Salvar
file_path = "processo_design_interfaces.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
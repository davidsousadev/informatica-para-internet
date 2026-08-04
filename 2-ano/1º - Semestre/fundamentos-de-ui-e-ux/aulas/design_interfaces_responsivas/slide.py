from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1 - Capa
add_slide(
    "Design de Interfaces Responsivas",
    "Princípios de Design Visual\nEnsino Médio"
)

# Slide 2 - O que é Design Responsivo
add_slide(
    "O que é Design Responsivo?",
    "É a prática de criar interfaces que se adaptam automaticamente "
    "ao tamanho da tela do dispositivo.\n\n"
    "Exemplos de dispositivos:\n"
    "- Smartphones\n"
    "- Tablets\n"
    "- Notebooks\n"
    "- Monitores\n"
    "- TVs"
)

# Slide 3 - Importância
add_slide(
    "Por que o Design Responsivo é Importante?",
    "- Melhora a experiência do usuário\n"
    "- Facilita a navegação\n"
    "- Mantém o conteúdo organizado\n"
    "- Funciona em diferentes tamanhos de tela\n"
    "- É essencial no desenvolvimento web moderno"
)

# Slide 4 - Características do Design Responsivo
add_slide(
    "Características do Design Responsivo",
    "- Layouts flexíveis\n"
    "- Uso de breakpoints\n"
    "- Sistemas de grade (grid system)\n"
    "- Fontes responsivas\n"
    "- Ajuste automático dos elementos"
)

# Slide 5 - Hierarquia Visual e Espaçamento
add_slide(
    "Hierarquia Visual e Espaçamento",
    "Os espaçamentos ajudam na organização e leitura da página.\n\n"
    "Tipos de espaçamento:\n"
    "- Padding: espaço interno\n"
    "- Margin: espaço externo\n\n"
    "Um bom espaçamento melhora a legibilidade."
)

# Slide 6 - Sistema de Grades
add_slide(
    "Sistema de Grades (Grid System)",
    "O grid system divide a interface em colunas.\n\n"
    "O modelo mais usado possui 12 colunas.\n\n"
    "Exemplo:\n"
    "- Metade da página = 6 colunas\n"
    "- Página inteira = 12 colunas\n"
    "- Menor elemento = 1 coluna"
)

# Slide 7 - Elementos do Grid
add_slide(
    "Elementos do Grid System",
    "O sistema de grades possui:\n\n"
    "- Margens\n"
    "- Canaletas (gutters)\n"
    "- Linhas (rows)\n"
    "- Colunas\n"
    "- Breakpoints"
)

# Slide 8 - Contêineres
add_slide(
    "Contêineres",
    "Os contêineres delimitam o espaço do conteúdo.\n\n"
    "Funções:\n"
    "- Organizar textos e imagens\n"
    "- Evitar conteúdo fora da tela\n"
    "- Melhorar a visualização em diferentes dispositivos"
)

# Slide 9 - Breakpoints
add_slide(
    "Breakpoints",
    "Breakpoints definem comportamentos diferentes para cada tamanho de tela.\n\n"
    "Exemplos no Bootstrap:\n"
    "- sm → smartphones\n"
    "- md → tablets\n"
    "- xxl → telas grandes"
)

# Slide 10 - Exemplo de Breakpoints
add_slide(
    "Exemplo de Uso dos Breakpoints",
    "Exemplo:\n\n"
    "- Tela com 500px usa a classe 'sm'\n"
    "- Telas acima de 1400px usam 'xxl'\n\n"
    "Isso permite adaptar automaticamente o layout."
)

# Slide 11 - Contêiner Fluido
add_slide(
    "Contêiner Fluido",
    "Um contêiner fluido ocupa 100% da largura disponível da tela.\n\n"
    "Vantagens:\n"
    "- Melhor aproveitamento do espaço\n"
    "- Adaptação automática\n"
    "- Maior flexibilidade"
)

# Slide 12 - Resumo
add_slide(
    "Resumo",
    "Aprendemos sobre:\n\n"
    "- Design responsivo\n"
    "- Espaçamento e hierarquia visual\n"
    "- Sistema de grades\n"
    "- Contêineres\n"
    "- Breakpoints\n"
    "- Layouts adaptáveis"
)

# Slide 13 - Atividade
add_activity_slide(
    "Atividade",
    "1. O que é design responsivo?\n\n"
    "2. Qual a função do grid system?\n\n"
    "3. Qual a diferença entre margin e padding?\n\n"
    "4. O que são breakpoints?\n\n"
    "5. Cite exemplos de dispositivos que usam design responsivo."
)

file_path = "design_interfaces_responsivas.pptx"
prs.save(file_path)

print(f"Apresentação salva em: {file_path}")

file_path
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Slide inicial
add_title_slide("Front-end vs. back-end e client-side vs. server-side", "Introdução ao UI/UX Design")

# Slides
add_slide("Conceitos Fundamentais", [
    "Ao desenvolver aplicações, você encontrará conceitos fundamentais",
    "Diferença entre front-end e back-end",
    "[ESPAÇO PARA EXPLICAÇÃO DO PROFESSOR]"
])

add_slide("Front-end", [
    "Parte visual com a qual os usuários interagem",
    "Elementos: botões, textos e imagens",
    "Uso de HTML, CSS e JavaScript",
    "[ESPAÇO PARA IMAGEM DE INTERFACE]"
])

add_slide("Profissionais de Front-end", [
    "Responsáveis por implementar o visual",
    "Podem colaborar com o design",
    "UI/UX geralmente define a experiência",
    "[ESPAÇO PARA EXEMPLOS]"
])

add_slide("Back-end", [
    "Parte que opera nos bastidores",
    "Processada em servidor remoto",
    "Gerencia dados e banco de dados",
    "[ESPAÇO PARA DIAGRAMA]"
])

add_slide("Novas Tecnologias", [
    "Front-end pode realizar parte do trabalho do servidor",
    "Envio de páginas já processadas em HTML",
    "Maior complexidade",
    "[ESPAÇO PARA EXPLICAÇÃO]"
])

add_slide("Client-side vs Server-side", [
    "Client-side: processado no navegador do usuário",
    "Server-side: processado no servidor remoto",
    "[ESPAÇO PARA COMPARAÇÃO VISUAL]"
])

add_slide("Exemplos Práticos", [
    "Animações → client-side",
    "Transações com cartão → server-side",
    "Segurança é essencial",
    "[INSERIR IMAGEM DA TRANSAÇÃO]"
])

add_slide("UI Design", [
    "Construção da parte visual da interface",
    "Foco na estética",
    "Elementos: botões, ícones, cores, tipografia",
    "[ESPAÇO PARA IMAGEM UI]"
])

add_slide("UX Design", [
    "Experiência do usuário ao usar a interface",
    "Facilidade de uso e navegação",
    "Objetivo: experiência agradável",
    "[ESPAÇO PARA EXEMPLOS]"
])

add_slide("Diferença entre UI e UX", [
    "UI → aparência",
    "UX → funcionalidade",
    "Áreas correlacionadas",
    "[ESPAÇO PARA ATIVIDADE]"
])

add_slide("Design Centrado no Usuário", [
    "Produzir aplicações que atendam necessidades reais",
    "Priorizar o usuário nas decisões",
    "Baseado em pesquisas",
    "[ESPAÇO PARA EXPLICAÇÃO]"
])

add_slide("Princípios do DCU", [
    "Resolver a raiz do problema",
    "Foco nas pessoas",
    "Abordagem sistêmica",
    "Testes rápidos e contínuos"
])

add_slide("Etapas do DCU", [
    "Análise",
    "Concepção",
    "Design",
    "Avaliação e repetição",
    "[ESPAÇO PARA FLUXOGRAMA]"
])

add_slide("Importância da Usabilidade", [
    "Interfaces intuitivas",
    "Priorizar facilidade de uso",
    "UX garante eficácia",
    "Um site bonito, mas difícil de usar não cumpre seu papel"
])

# Salvar
file_path = "front-endxback-end.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
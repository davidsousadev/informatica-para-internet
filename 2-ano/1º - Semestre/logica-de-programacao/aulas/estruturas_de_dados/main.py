from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1
add_slide(
    "Estruturas de Dados em Portugol",
    "Conceitos fundamentais de vetores, matrizes, registros, pilhas, filas e dicionários."
)

# Slide 2
add_slide(
    "Objetivos da Aula",
    "- Compreender estruturas de dados.\n"
    "- Utilizar vetores e matrizes.\n"
    "- Conhecer registros.\n"
    "- Entender pilhas e filas.\n"
    "- Aplicar dicionários/mapas."
)

# Slide 3
add_slide(
    "O que são Estruturas de Dados?",
    "São formas de organizar e armazenar informações na memória para facilitar o acesso e manipulação dos dados."
)

# Slide 4
add_slide(
    "Tipos de Dados Primitivos",
    "Inteiro, Real, Caractere, Cadeia e Lógico.\n"
    "Exemplos:\n"
    "idade <- 20\n"
    "nota <- 8.5\n"
    "nome <- \"Maria\"\n"
    "aprovado <- verdadeiro"
)

# Slide 5
add_slide(
    "Vetores",
    "Vetores armazenam vários valores do mesmo tipo em uma única variável.\n\n"
    "Exemplo:\n"
    "vetor notas[5]"
)

# Slide 6
add_slide(
    "Exemplo de Vetor em Portugol",
    "algoritmo \"vetor\"\n"
    "var\n"
    " notas: vetor[1..5] de real\n"
    "inicio\n"
    " notas[1] <- 7.5\n"
    " escreva(notas[1])\n"
    "fimalgoritmo"
)

# Slide 7
add_slide(
    "Percorrendo Vetores",
    "algoritmo \"percorrer\"\n"
    "var\n"
    " i: inteiro\n"
    " numeros: vetor[1..5] de inteiro\n"
    "inicio\n"
    " para i de 1 ate 5 faca\n"
    "   escreva(numeros[i])\n"
    " fimpara\n"
    "fimalgoritmo"
)

# Slide 8
add_slide(
    "Matrizes",
    "Matrizes são estruturas bidimensionais.\n"
    "Possuem linhas e colunas.\n"
    "Exemplo: matriz[1..3,1..3]"
)

# Slide 9
add_slide(
    "Exemplo de Matriz",
    "algoritmo \"matriz\"\n"
    "var\n"
    " tabela: matriz[1..3,1..3] de inteiro\n"
    "inicio\n"
    " tabela[1,1] <- 10\n"
    " escreva(tabela[1,1])\n"
    "fimalgoritmo"
)

# Slide 10
add_slide(
    "Registros",
    "Permitem armazenar diferentes tipos de informações relacionadas em uma única estrutura."
)

# Slide 11
add_slide(
    "Exemplo de Registro",
    "tipo Pessoa = registro\n"
    " nome: cadeia\n"
    " idade: inteiro\n"
    " fimregistro\n\n"
    "var aluno: Pessoa"
)

# Slide 12
add_slide(
    "Pilhas (LIFO)",
    "Last In, First Out.\n"
    "O último elemento inserido é o primeiro a sair.\n\n"
    "Exemplo: pilha de pratos."
)

# Slide 13
add_slide(
    "Operações em Pilhas",
    "Empilhar (Push)\n"
    "Desempilhar (Pop)\n"
    "Topo da Pilha"
)

# Slide 14
add_slide(
    "Filas (FIFO)",
    "First In, First Out.\n"
    "O primeiro elemento inserido é o primeiro a sair.\n\n"
    "Exemplo: fila de banco."
)

# Slide 15
add_slide(
    "Operações em Filas",
    "Enfileirar\n"
    "Desenfileirar\n"
    "Início e Fim da Fila"
)

# Slide 16
add_slide(
    "Dicionários e Mapas",
    "Estruturas compostas por chave e valor.\n"
    "Permitem localizar informações rapidamente."
)

# Slide 17
add_slide(
    "Métodos e Estruturas",
    "Métodos são funções associadas a estruturas ou objetos.\n"
    "Exemplos:\n"
    "- Inserir elemento\n"
    "- Remover elemento\n"
    "- Buscar elemento\n"
    "- Ordenar dados"
)

# Slide 18
add_slide(
    "Recursividade",
    "Uma função que chama a si mesma.\n"
    "Necessita de um caso base para evitar execução infinita."
)

# Slide 19
add_slide(
    "Exemplo Conceitual de Recursão",
    "Fatorial:\n"
    "5! = 5 * 4!\n"
    "4! = 4 * 3!\n"
    "...\n"
    "1! = 1 (caso base)"
)

# Slide 20
add_slide(
    "Resumo",
    "- Vetores: uma dimensão.\n"
    "- Matrizes: duas dimensões.\n"
    "- Registros: dados agrupados.\n"
    "- Pilhas: LIFO.\n"
    "- Filas: FIFO.\n"
    "- Dicionários: chave e valor."
)

# Atividade
add_activity_slide(
    "Atividade Prática",
    "1. Crie um vetor com 5 números e exiba todos.\n\n"
    "2. Crie uma matriz 3x3 e preencha com valores.\n\n"
    "3. Crie um registro Aluno contendo nome e idade.\n\n"
    "4. Explique a diferença entre Pilha e Fila.\n\n"
    "5. Cite uma aplicação prática para dicionários."
)

file_path = "aula_07_estruturas_de_Dados_portugol.pptx"
prs.save(file_path)

print(file_path)
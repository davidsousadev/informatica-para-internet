from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


# Slide inicial
add_title_slide("Operações Aritméticas em Portugol", "Lógica de Programação")

# Objetivo
add_slide("Objetivo da Aula", [
    "Compreender operações matemáticas",
    "Utilizar cálculos em programas",
    "Aplicar operações com variáveis"
])

# Conceitos
add_slide("Conceitos Fundamentais", [
    "Soma → +",
    "Subtração → -",
    "Multiplicação → *",
    "Divisão → /"
])

# Soma
add_slide("Soma", [
    "Utilizada para adicionar valores",
    "Exemplo:",
    "soma = num1 + num2"
])

# Subtração
add_slide("Subtração", [
    "Utilizada para calcular diferença",
    "Exemplo:",
    "sub = num1 - num2"
])

# Multiplicação
add_slide("Multiplicação", [
    "Utilizada para multiplicar valores",
    "Exemplo:",
    "mult = num1 * num2"
])

# Divisão
add_slide("Divisão", [
    "Utilizada para dividir valores",
    "Exemplo:",
    "div = num1 / num2"
])

# Exemplo completo
add_slide("Exemplo Completo", [
    "programa {",
    "funcao inicio() {",
    "real num1, num2, soma",
    "leia(num1)",
    "leia(num2)",
    "soma = num1 + num2",
    "escreva(\"Resultado: \", soma)",
    "}",
    "}"
])

# Atividade 1
add_slide("Atividade 1 - Soma com Nome", [
    "Arquivo: soma_nome.por",
    "Ler nome e dois números",
    "Exibir nome + resultado da soma"
])

# Atividade 2
add_slide("Atividade 2 - Idade + Multiplicação", [
    "Arquivo: idade_multiplicacao.por",
    "Ler idade e dois números",
    "Exibir idade + multiplicação"
])

# Atividade 3
add_slide("Atividade 3 - Produto com Cálculo", [
    "Arquivo: produto_calculo.por",
    "Ler produto, valor e quantidade",
    "Calcular total da compra"
])

# Atividade 4
add_slide("Atividade 4 - Subtração", [
    "Arquivo: boas_subtracao.por",
    "Ler nome e dois números",
    "Exibir boas-vindas + subtração"
])

# Atividade 5
add_slide("Atividade 5 - Divisão", [
    "Arquivo: disciplina_divisao.por",
    "Ler nome, disciplina e números",
    "Exibir dados + divisão"
])

# Encerramento
add_slide("Resumo", [
    "Uso de + - * /",
    "Operações com variáveis",
    "Integração com leia e escreva",
    "Base para lógica de programação"
])

# Desafio
add_slide("Desafio", [
    "Criar um programa que:",
    "Leia dois números",
    "Mostre soma, subtração, multiplicação e divisão"
])


# Salvar
file_path = "aula_04_operacoes_aritimeticas_portugol.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
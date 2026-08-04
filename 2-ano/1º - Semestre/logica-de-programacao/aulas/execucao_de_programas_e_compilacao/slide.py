from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1 - Capa

add_slide(
"Execução de Programas e Compilação",
"Como os programas são executados?\nLinguagens Interpretadas e Compiladas\nEnsino Médio"
)

# Slide 2 - Introdução

add_slide(
"O que é um Programa?",
"Um programa é um conjunto de instruções que o computador executa para realizar uma tarefa.\n\nEssas instruções são escritas em linguagens de programação."
)

# Slide 3 - Linguagens Interpretadas

add_slide(
"Linguagens Interpretadas",
"São linguagens que precisam de um interpretador para executar o código.\n\nExemplos: Python, JavaScript e Ruby."
)

# Slide 4 - Criando um Script Python

add_slide(
"Etapa 1 - Criando o script.py",
"No VS Code, crie um arquivo chamado script.py.\n\nExemplo:\nprint('Olá, mundo!')\nprint('Bem-vindo ao Python!')"
)

# Slide 5 - Executando o Script

add_slide(
"Etapa 2 - Executando o Script",
"Abra o terminal e execute:\n\npython script.py\n\nO interpretador Python lerá o código e mostrará o resultado na tela."
)

# Slide 6 - O que é o Interpretador?

add_slide(
"Interpretador Python",
"O interpretador traduz e executa cada linha do código em tempo real.\n\nÉ necessário instalar o Python para executar scripts."
)

# Slide 7 - VS Code

add_slide(
"Executando no VS Code",
"O VS Code possui um terminal integrado.\n\nMenu: Terminal > Novo Terminal\n\nTambém é possível executar o código diretamente pelo editor."
)

# Slide 8 - O que é Compilação?

add_slide(
"Compilação",
"Em linguagens compiladas, o código-fonte é convertido para linguagem de máquina antes da execução.\n\nO resultado é um arquivo executável."
)

# Slide 9 - Processo de Compilação

add_slide(
"Fluxo da Compilação",
"Código-Fonte\n        ↓\nCompilador\n        ↓\nArquivo Executável"
)

# Slide 10 - Vantagens e Desvantagens

add_slide(
"Vantagens e Desvantagens",
"Vantagens:\n• Maior desempenho\n• Execução mais rápida\n\nDesvantagens:\n• É necessário compilar\n• Executáveis podem ser maiores"
)

# Slide 11 - Exemplo em C++

add_slide(
"Exemplo de Código C++",
"#include <iostream>\n\nint main() {\n  std::cout << \"Ola mundo!\";\n  return 0;\n}"
)

# Slide 12 - Compilando o Código

add_slide(
"Compilando com g++",
"Comando:\n\ng++ codigo.cpp -o programa_executavel\n\nO compilador gera um arquivo executável."
)

# Slide 13 - Executando o Programa

add_slide(
"Executando o Programa",
"Windows:\nClique duas vezes no executável.\n\nLinux/macOS:\n./programa_executavel"
)

# Slide 14 - Comparação

add_slide(
"Interpretado x Compilado",
"Interpretado (Python):\nCódigo → Interpretador → Execução\n\nCompilado (C++):\nCódigo → Compilador → Executável → Execução"
)

# Slide 15 - Resumo

add_slide(
"Resumo da Aula",
"• Python é uma linguagem interpretada.\n• C++ é uma linguagem compilada.\n• O interpretador executa o código diretamente.\n• O compilador gera um executável."
)

# Slide 16 - Atividade

add_slide(
"Atividade",
"1. O que faz um interpretador?\n\n2. Qual a diferença entre linguagem interpretada e compilada?\n\n3. Cite um exemplo de cada tipo."
)

# Slide 17 - Desafio

add_slide(
"Desafio",
"Crie um arquivo script.py que exiba três mensagens diferentes no terminal.\n\nExecute utilizando:\npython script.py"
)

file_path = "Execucao_Programas_e_Compilacao.pptx"
prs.save(file_path)

print(file_path)
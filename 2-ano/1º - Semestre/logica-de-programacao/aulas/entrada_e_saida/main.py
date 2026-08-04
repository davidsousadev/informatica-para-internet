from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    
    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


# Slide inicial
add_title_slide("Entrada e Saída de Dados em Portugol", "Lógica de Programação")

# Objetivo
add_slide("Objetivo da Aula", [
    "Compreender entrada e saída de dados",
    "Utilizar leia e escreva",
    "Criar programas simples com variáveis"
])

# Conceitos
add_slide("Conceitos Fundamentais", [
    "leia → entrada de dados",
    "escreva → saída de dados",
    "Programa interage com o usuário"
])

# Estrutura
add_slide("Estrutura Básica do Portugol", [
    "algoritmo \"nome\"",
    "var → declaração de variáveis",
    "inicio → início do programa",
    "fimalgoritmo → fim do programa"
])

# Exemplo
add_slide("Exemplo Simples", [
    "algoritmo \"exemplo\"",
    "var nome: cadeia",
    "inicio",
    "escreva(\"Digite seu nome: \")",
    "leia(nome)",
    "escreva(\"Olá, \", nome)",
    "fimalgoritmo"
])

# Atividade 1
add_slide("Atividade 1 - Mostrar Idade", [
    "Ler a idade do usuário",
    "Exibir a idade informada",
    "[IMPLEMENTAR NO PORTUGOL]"
])

add_slide("Código - Mostrar Idade", [
    "algoritmo \"mostrar_idade\"",
    "var idade: inteiro",
    "inicio",
    "leia(idade)",
    "escreva(\"Você tem \", idade, \" anos\")",
    "fimalgoritmo"
])

# Atividade 2
add_slide("Atividade 2 - Cadastro Simples", [
    "Ler nome e cidade",
    "Exibir frase com os dados",
    "[IMPLEMENTAR NO PORTUGOL]"
])

add_slide("Código - Cadastro Simples", [
    "algoritmo \"cadastro_simples\"",
    "var nome, cidade: cadeia",
    "inicio",
    "leia(nome)",
    "leia(cidade)",
    "escreva(\"Seu nome é \", nome, \" e mora em \", cidade)",
    "fimalgoritmo"
])

# Atividade 3
add_slide("Atividade 3 - Produto e Quantidade", [
    "Ler nome do produto",
    "Ler quantidade",
    "Exibir informações"
])

add_slide("Código - Produto", [
    "algoritmo \"produto_quantidade\"",
    "var produto: cadeia",
    "var quantidade: inteiro",
    "inicio",
    "leia(produto)",
    "leia(quantidade)",
    "escreva(\"Produto: \", produto)",
    "escreva(\"Quantidade: \", quantidade)",
    "fimalgoritmo"
])

# Atividade 4
add_slide("Atividade 4 - Boas-vindas", [
    "Ler nome",
    "Exibir mensagem personalizada"
])

add_slide("Código - Boas-vindas", [
    "algoritmo \"boas_vindas\"",
    "var nome: cadeia",
    "inicio",
    "leia(nome)",
    "escreva(\"Bem-vindo, \", nome)",
    "fimalgoritmo"
])

# Atividade 5
add_slide("Atividade 5 - Disciplina Favorita", [
    "Ler nome",
    "Ler disciplina favorita",
    "Exibir frase"
])

add_slide("Código - Disciplina", [
    "algoritmo \"disciplina_favorita\"",
    "var nome, disciplina: cadeia",
    "inicio",
    "leia(nome)",
    "leia(disciplina)",
    "escreva(nome, \" gosta de \", disciplina)",
    "fimalgoritmo"
])

# Encerramento
add_slide("Resumo", [
    "Uso de variáveis",
    "Entrada com leia",
    "Saída com escreva",
    "Construção de frases"
])

# Desafio
add_slide("Desafio", [
    "Criar um programa que leia:",
    "nome, idade e cidade",
    "Exibir tudo em uma única frase"
])


# Salvar
file_path = "aula_02_entrada_e_saida_portugol.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
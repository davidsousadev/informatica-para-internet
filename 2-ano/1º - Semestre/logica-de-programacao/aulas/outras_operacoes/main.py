from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


# Slide Inicial
add_title_slide(
    "Outras Operações em Portugol",
    "Lógica de Programação"
)

# Objetivo
add_slide("Objetivo da Aula", [
    "Compreender operações matemáticas",
    "Utilizar biblioteca Matemática",
    "Aplicar potência e raiz quadrada",
    "Desenvolver cálculos em Portugol"
])

# Biblioteca Matemática
add_slide("Biblioteca Matemática", [
    "A biblioteca Matemática possui funções prontas",
    "Facilita cálculos matemáticos",
    "Precisamos incluir a biblioteca no programa",
    "Exemplo:",
    "inclua biblioteca Matematica --> mat"
])

# Prioridade das Operações
add_slide("Prioridade das Operações", [
    "Os cálculos seguem uma ordem matemática",
    "Parênteses possuem prioridade",
    "Multiplicação e divisão vêm antes da soma",
    "Exemplo:",
    "8 + 2 * 3 = 14",
    "(8 + 2) * 3 = 30"
])

# Operador de resto
add_slide("Operador de Resto (%)", [
    "O operador % mostra o resto da divisão",
    "Muito usado em lógica de programação",
    "Exemplos:",
    "15 % 4 = 3",
    "17 % 5 = 2"
])

# Potência
add_slide("Função Potência", [
    "Usamos mat.potencia()",
    "Calcula um número elevado a outro",
    "Exemplo:",
    "mat.potencia(valor, 3.0)",
    "Calcula o número elevado ao cubo"
])

# Raiz Quadrada
add_slide("Função Raiz", [
    "Usamos mat.raiz()",
    "Calcula a raiz de um número",
    "Exemplo:",
    "mat.raiz(valor, 2.0)",
    "Calcula a raiz quadrada"
])

# Explicação do Código
add_slide("Explicação do Código", [
    "leia(valor) → recebe o número digitado",
    "mat.potencia(valor, 3.0) → cubo do número",
    "mat.raiz(valor, 2.0) → raiz quadrada",
    "escreva() → mostra os resultados"
])

# Código Completo
add_slide("Exemplo Completo em Portugol", [
    "programa {",
    "inclua biblioteca Matematica --> mat",
    "funcao inicio() {",
    "real valor, potencia, raiz_quadrada",
    "leia(valor)",
    "potencia = mat.potencia(valor, 3.0)",
    "raiz_quadrada = mat.raiz(valor, 2.0)",
    "escreva(potencia)",
    "escreva(raiz_quadrada)",
    "}",
    "}"
])

# Aplicações práticas
add_slide("Aplicações Práticas", [
    "Cálculos matemáticos",
    "Sistemas financeiros",
    "Jogos digitais",
    "Aplicativos de engenharia",
    "Programação científica"
])

# Atividades da Aula
add_slide("Atividades Propostas", [
    "Prioridade das operações",
    "Metade inteira e resto",
    "Potência e raiz quadrada",
    "Número ao cubo",
    "Média com prioridade"
])

# Resumo
add_slide("Resumo da Aula", [
    "Uso da biblioteca Matemática",
    "Aplicação de potência",
    "Uso de raiz quadrada",
    "Operador de resto (%)",
    "Cálculos em Portugol"
])

# Desafio
add_slide("Desafio Final", [
    "Criar um programa que:",
    "Leia um número",
    "Mostre quadrado, cubo e raiz",
    "Mostre o resto da divisão por 2"
])

# Salvar apresentação
file_path = "aula_05_outras_operacoes.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)
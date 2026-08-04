from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1 - Capa
add_slide(
    "Realidade Aumentada e Virtual (RA/RV): Conteúdo em Qualquer Lugar!",
    "Tecnologias imersivas que transformam a forma como interagimos com o mundo digital."
)

# Slide 2 - Introdução
add_slide(
    "O que são RA e RV?",
    "A Realidade Aumentada (RA) adiciona elementos digitais ao mundo real. "
    "A Realidade Virtual (RV) cria ambientes totalmente digitais e imersivos. "
    "Ambas ampliam as possibilidades de interação e aprendizado."
)

# Slide 3 - Diferença entre RA e RV
add_slide(
    "RA x RV",
    "RA: mistura elementos virtuais ao ambiente físico.\n"
    "RV: substitui completamente o ambiente real por um virtual.\n"
    "Exemplo RA: visualizar móveis em casa pelo celular.\n"
    "Exemplo RV: explorar um ambiente virtual usando óculos VR."
)

# Slide 4 - Realidade Mista
add_slide(
    "Realidade Mista (MR)",
    "Combina recursos de RA e RV.\n"
    "Permite interação simultânea entre objetos físicos e virtuais.\n"
    "Muito utilizada em treinamentos avançados e aplicações industriais."
)

# Slide 5 - Principais Tipos de Óculos
add_slide(
    "Tipos de Óculos RA/RV",
    "• Óculos de Realidade Aumentada (RA)\n"
    "• Óculos de Realidade Virtual (RV)\n"
    "• Óculos de Realidade Mista (MR)\n"
    "Cada modelo atende necessidades específicas de uso."
)

# Slide 6 - Características Importantes
add_slide(
    "Características dos Dispositivos",
    "• Campo de Visão (FOV)\n"
    "• Resolução e qualidade de imagem\n"
    "• Sensores de movimento\n"
    "• Rastreamento ocular\n"
    "• Portabilidade e conectividade"
)

# Slide 7 - Aplicações na Educação
add_slide(
    "Educação e Treinamento",
    "Permitem simulações práticas em ambientes seguros.\n"
    "Facilitam o aprendizado interativo.\n"
    "Utilizadas em medicina, engenharia, aviação e áreas técnicas."
)

# Slide 8 - Aplicações na Saúde
add_slide(
    "RA/RV na Saúde",
    "• Treinamento de profissionais.\n"
    "• Reabilitação física.\n"
    "• Tratamento de fobias e ansiedade.\n"
    "• Simulações cirúrgicas com maior precisão."
)

# Slide 9 - Aplicações em Arquitetura e Design
add_slide(
    "Arquitetura e Design",
    "Visualização de projetos em tamanho real.\n"
    "Identificação de melhorias antes da construção.\n"
    "Redução de custos e aumento da eficiência."
)

# Slide 10 - Jogos e Entretenimento
add_slide(
    "Jogos e Entretenimento",
    "Experiências mais imersivas e interativas.\n"
    "Usuários participam ativamente dos cenários virtuais.\n"
    "Uma das áreas que mais impulsionam a evolução da tecnologia."
)

# Slide 11 - Benefícios da RA/RV
add_slide(
    "Principais Benefícios",
    "• Maior imersão.\n"
    "• Aprendizagem prática.\n"
    "• Redução de riscos em treinamentos.\n"
    "• Aumento da produtividade.\n"
    "• Novas formas de interação digital."
)

# Slide 12 - Curiosidade
add_slide(
    "Curiosidade",
    "A NASA utiliza Realidade Virtual para treinar astronautas, "
    "permitindo simular missões espaciais antes das operações reais."
)

# Slide 13 - Desafios
add_slide(
    "Desafios da Tecnologia",
    "• Alto custo de alguns equipamentos.\n"
    "• Necessidade de hardware especializado.\n"
    "• Conforto do usuário em longos períodos.\n"
    "• Desenvolvimento de conteúdo de qualidade."
)

# Slide 14 - Futuro da RA/RV
add_slide(
    "O Futuro da RA/RV",
    "Integração com Inteligência Artificial.\n"
    "Maior uso no trabalho remoto.\n"
    "Expansão para educação, saúde, indústria e comércio.\n"
    "Dispositivos cada vez mais acessíveis."
)

# Slide 15 - Atividade
add_activity_slide(
    "Atividade",
    "1. Qual a principal diferença entre RA e RV?\n"
    "2. Cite duas aplicações da RA/RV na saúde.\n"
    "3. Como a RA/RV pode melhorar o aprendizado?\n"
    "4. Qual aplicação você considera mais interessante? Justifique."
)

# Slide 16 - Encerramento
add_slide(
    "Conclusão",
    "A Realidade Aumentada e a Realidade Virtual estão transformando a forma "
    "como aprendemos, trabalhamos e nos divertimos, tornando o conteúdo acessível em qualquer lugar."
)

file_path = "Realidade_Aumentada_e_Virtual_RA_RV.pptx"
prs.save(file_path)

print(file_path)
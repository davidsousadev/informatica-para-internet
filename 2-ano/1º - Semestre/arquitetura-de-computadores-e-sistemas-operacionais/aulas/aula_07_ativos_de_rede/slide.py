from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1 - Capa
add_slide(
    "Aula 07 - Ativos de Rede",
    "Roteadores, Switches e Servidores\n\n"
    "Redes de Comunicação e Arquitetura Cliente/Servidor"
)

# Slide 2 - Contextualização
add_slide(
    "Contextualizando",
    "Vivemos em um mundo hiperconectado onde as redes de comunicação "
    "permitem a troca de informações entre dispositivos.\n\n"
    "Ao acessar sites, enviar mensagens ou assistir vídeos, utilizamos "
    "redes que conectam sistemas e pessoas."
)

# Slide 3 - Curiosidade
add_slide(
    "Curiosidade",
    "A primeira mensagem da Internet foi enviada em 1969.\n\n"
    "Desde então, as redes evoluíram até possibilitar a Internet moderna, "
    "conectando bilhões de dispositivos em todo o mundo."
)

# Slide 4 - O que são Redes de Comunicação
add_slide(
    "Redes de Comunicação",
    "As redes permitem que dispositivos interajam e compartilhem informações.\n\n"
    "Podemos imaginar cada dispositivo como uma ilha e a rede como uma ponte "
    "que conecta essas ilhas para a troca de dados."
)

# Slide 5 - Arquitetura Cliente/Servidor
add_slide(
    "Arquitetura Cliente/Servidor",
    "Modelo mais utilizado nas redes modernas.\n\n"
    "Cliente: solicita informações ou serviços.\n"
    "Servidor: processa solicitações e envia respostas.\n\n"
    "Exemplo: navegador acessando um site."
)

# Slide 6 - Cliente
add_slide(
    "O Cliente",
    "É o dispositivo que solicita informações ou serviços.\n\n"
    "Exemplos:\n"
    "• Computador\n"
    "• Notebook\n"
    "• Smartphone\n"
    "• Tablet\n\n"
    "Quando você acessa um site, seu dispositivo atua como cliente."
)

# Slide 7 - Servidor
add_slide(
    "O Servidor",
    "Recebe solicitações dos clientes, processa dados e envia respostas.\n\n"
    "Pode atender centenas ou milhares de usuários simultaneamente."
)

# Slide 8 - O que são Ativos de Rede
add_slide(
    "Ativos de Rede",
    "São dispositivos físicos e lógicos que formam a infraestrutura da rede.\n\n"
    "Principais exemplos:\n"
    "• Roteadores\n"
    "• Switches\n"
    "• Servidores\n"
    "• Access Points\n"
    "• Firewalls"
)

# Slide 9 - Roteadores
add_slide(
    "Roteadores",
    "Conectam redes diferentes entre si.\n\n"
    "Função principal:\n"
    "• Encaminhar pacotes de dados.\n"
    "• Interligar a rede local à Internet.\n\n"
    "Exemplo: roteador residencial."
)

# Slide 10 - Características dos Roteadores
add_slide(
    "Características dos Roteadores",
    "• Possuem endereços IP internos e externos.\n"
    "• Escolhem a melhor rota para os dados.\n"
    "• Muitos modelos incluem Wi-Fi e switch integrado.\n\n"
    "Curiosidade: a maioria possui firewall básico."
)

# Slide 11 - Switches
add_slide(
    "Switches",
    "Conectam dispositivos dentro da mesma rede local (LAN).\n\n"
    "Exemplos de dispositivos conectados:\n"
    "• Computadores\n"
    "• Impressoras\n"
    "• Câmeras\n"
    "• Servidores"
)

# Slide 12 - Características dos Switches
add_slide(
    "Características dos Switches",
    "• Organizam o tráfego de dados.\n"
    "• Criam comunicação direta entre dispositivos.\n"
    "• Reduzem colisões na rede.\n\n"
    "Muito utilizados em empresas."
)

# Slide 13 - Servidores
add_slide(
    "Servidores",
    "São responsáveis por armazenar dados e fornecer serviços aos clientes.\n\n"
    "Funcionam continuamente para garantir disponibilidade das informações."
)

# Slide 14 - Tipos de Servidores
add_slide(
    "Tipos de Servidores",
    "• Servidor Web\n"
    "• Servidor de Arquivos\n"
    "• Servidor de Banco de Dados\n"
    "• Servidor de E-mail\n\n"
    "Cada um possui funções específicas na rede."
)

# Slide 15 - Access Points
add_slide(
    "Access Points",
    "Dispositivos que ampliam a cobertura da rede sem fio.\n\n"
    "Permitem acesso Wi-Fi em locais onde o sinal do roteador não alcança."
)

# Slide 16 - Firewalls
add_slide(
    "Firewalls",
    "Funcionam como uma barreira de proteção.\n\n"
    "Monitoram e controlam o tráfego de entrada e saída da rede.\n"
    "Bloqueiam acessos não autorizados e ameaças."
)

# Slide 17 - Comparação dos Ativos
add_slide(
    "Comparação dos Ativos",
    "Roteador: conecta redes diferentes.\n\n"
    "Switch: conecta dispositivos da mesma rede.\n\n"
    "Servidor: fornece serviços e dados aos clientes.\n\n"
    "Access Point: amplia o Wi-Fi.\n\n"
    "Firewall: protege a rede."
)

# Slide 18 - Aplicações no Dia a Dia
add_slide(
    "Aplicações no Dia a Dia",
    "Ao usar redes sociais, streaming, jogos online ou e-mail:\n\n"
    "• Cliente solicita informações.\n"
    "• Servidor responde.\n"
    "• Roteadores e switches transportam os dados.\n"
    "• Firewalls ajudam na segurança."
)

# Slide 19 - Resumo da Aula
add_slide(
    "Resumo",
    "✓ Redes conectam dispositivos.\n"
    "✓ Cliente solicita serviços.\n"
    "✓ Servidor fornece serviços.\n"
    "✓ Roteadores conectam redes.\n"
    "✓ Switches conectam dispositivos locais.\n"
    "✓ Firewalls aumentam a segurança."
)

# Slide 20 - Atividade
add_activity_slide(
    "Atividade",
    "1. Qual a diferença entre cliente e servidor?\n\n"
    "2. Qual é a função principal de um roteador?\n\n"
    "3. Para que serve um switch?\n\n"
    "4. Cite dois tipos de servidores.\n\n"
    "5. Como o firewall contribui para a segurança da rede?"
)

file_path = "Aula_07_Ativos_de_Rede.pptx"
prs.save(file_path)

print(file_path)
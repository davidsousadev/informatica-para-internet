from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slides
add_slide(
    "Vamos começar com o básico",
    "Os sistemas que fazem o hardware funcionar!"
)

add_slide(
    "1.1.1 Firmware - Conceito",
    "Firmware é como o “DNA ou software embutido” em qualquer dispositivo, mantendo a conexão entre hardware e sistema operacional.\n\n"
    "Assim que você liga seu dispositivo, é o firmware que entra em ação, preparando o terreno para tudo o que virá depois."
)

add_slide(
    "1.1.1 Firmware - Curiosidade",
    "🔍 CURIOSIDADE\n\nSabia que alguns aparelhos, como roteadores e impressoras, raramente precisam atualizar o firmware?\n\n"
    "Já dispositivos móveis recebem atualizações regulares para melhorias de segurança e desempenho!"
)

add_slide(
    "1.1.1 Firmware - Exemplo Clássico",
    "💡 SAIBA MAIS\n\nA BIOS, aquela tela preta que aparece quando você liga o computador, é um tipo clássico de firmware.\n\n"
    "Ela é fundamental para iniciar todo o funcionamento da máquina antes do sistema operacional carregar."
)

add_slide(
    "1.1.2 Sistemas Operacionais - Definição",
    "O Sistema Operacional (SO) é o grande maestro do seu dispositivo.\n\n"
    "Ele organiza os recursos e permite que você interaja com o computador ou celular de forma prática e intuitiva."
)

add_slide(
    "1.1.2 Sistemas Operacionais - Exemplos",
    "Existem vários tipos de sistemas operacionais, cada um com suas características e funções específicas:\n\n"
    "- Windows\n"
    "- macOS\n"
    "- Linux\n"
    "- Android\n"
    "- iOS"
)

add_slide(
    "1.1.2 Sistemas Operacionais - História e Evolução",
    "Figura 2 – Windows 95 marcou época por facilitar o uso de computadores pessoais\n\n"
    "Disponível em: https://www.shutterstock.com. Acesso em: 21 nov 2024.\n\n"
    "O lançamento desse sistema foi um marco para popularizar o uso de computadores pessoais no mundo todo."
)

add_slide(
    "1.1.2 Sistemas Operacionais - Curiosidade",
    "🔍 CURIOSIDADE\n\nO Android, sistema operacional mais popular do mundo, roda em mais de 2 bilhões de dispositivos ao redor do planeta!\n\n"
    "Sua flexibilidade e variedade de aparelhos compatíveis explicam seu grande alcance."
)

add_slide(
    "1.1.2 Drivers - O que são?",
    "Os drivers são como “intérpretes” que traduzem as informações entre o sistema operacional e os componentes de hardware.\n\n"
    "Eles garantem que tudo funcione harmoniosamente, permitindo que cada parte do equipamento se comunique corretamente com o resto do sistema."
)

add_slide(
    "1.1.2 Drivers - Função e Importância",
    "Pense nos drivers como os responsáveis por permitir que o seu teclado, impressora, webcam, placa de vídeo e outros componentes operem sem problemas.\n\n"
    "Sem o driver correto, o hardware pode não funcionar ou ter seu funcionamento comprometido."
)

add_slide(
    "1.1.2 Drivers - Dicas Importantes",
    "✅ DICAS\n\nSempre mantenha os drivers atualizados!\n\n"
    "Atualizações frequentes podem melhorar muito a performance, corrigir erros e aumentar a segurança do seu sistema contra falhas ou problemas de compatibilidade."
)

file_path = "Aula 05 - Sistemas Operacionais, Firmware e Drivers.pptx"
prs.save(file_path)
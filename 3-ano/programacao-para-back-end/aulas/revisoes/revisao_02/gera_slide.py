from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(title, bullets):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)

def add_code_slide(title, code):
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title

    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(3))
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    p.text = code
    p.font.name = "Courier New"
    p.font.size = Pt(18)

# Capa
add_title_slide(
    "Variáveis em Python",
    "Disciplina: Programação para Back-End"
)

# Slides de conteúdo
add_bullets_slide(
    "O que são variáveis?",
    [
        "Variáveis armazenam dados na memória.",
        "São utilizadas para guardar informações temporárias.",
        "Em Python, criamos variáveis usando o sinal =.",
        "Python identifica automaticamente o tipo do dado."
    ]
)

add_code_slide(
    "Criando Variáveis",
    '''nome = "Maria"
idade = 25
altura = 1.70
ativo = True

print(nome)
print(idade)'''
)

add_bullets_slide(
    "Tipos de Dados",
    [
        "str → textos",
        "int → números inteiros",
        "float → números decimais",
        "bool → verdadeiro ou falso"
    ]
)

add_code_slide(
    "Exemplo de Tipos",
    '''cidade = "Fortaleza"
ano = 2026
temperatura = 29.5
online = True'''
)

add_bullets_slide(
    "Regras para Variáveis",
    [
        "Não podem começar com números.",
        "Não podem conter espaços.",
        "Use nomes claros e objetivos.",
        "Snake_case é o padrão do Python."
    ]
)

add_code_slide(
    "Boas Práticas",
    '''# Correto
nome_completo = "João"

# Evite
nc = "João"'''
)

add_bullets_slide(
    "Entrada de Dados",
    [
        "A função input() recebe dados do usuário.",
        "Os dados recebidos são textos.",
        "Podemos converter usando int() e float()."
    ]
)

add_code_slide(
    "Exemplo com input()",
    '''nome = input("Digite seu nome: ")
idade = int(input("Digite sua idade: "))

print(nome)
print(idade)'''
)

add_bullets_slide(
    "Resumo",
    [
        "Variáveis são essenciais na programação.",
        "Python possui tipagem dinâmica.",
        "Boas práticas melhoram a leitura do código.",
        "Variáveis facilitam o desenvolvimento Back-End."
    ]
)

add_bullets_slide(
    "Atividade Prática",
    [
        "1. Crie uma variável chamada nome.",
        "2. Crie uma variável chamada idade.",
        "3. Exiba os valores usando print().",
        "4. Faça um programa pedindo dados ao usuário."
    ]
)

# Salvar apresentação
file_path = "variaveis_python_backend.pptx"
prs.save(file_path)

print(f"Apresentação salva em: {file_path}")

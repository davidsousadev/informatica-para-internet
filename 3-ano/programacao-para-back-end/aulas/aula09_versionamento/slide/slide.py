from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1

add_slide(
"Versionamento de Código e APIs",
"""Como equipes de desenvolvimento controlam mudanças
em sistemas modernos sem perder informações."""
)

# Slide 2

add_slide(
"O Problema",
"""Imagine um trabalho em grupo:

• João altera um arquivo.
• Maria altera o mesmo arquivo.
• Pedro apaga algo sem querer.

Como recuperar versões anteriores?

Solução: Versionamento."""
)

# Slide 3

add_slide(
"O que é Versionamento?",
"""Versionamento é o processo de registrar mudanças
realizadas em um software ao longo do tempo.

Objetivos:
• Histórico de alterações
• Recuperação de versões
• Trabalho colaborativo
• Controle de qualidade"""
)

# Slide 4

add_slide(
"Git: A Ferramenta Mais Utilizada",
"""Git é um Sistema de Controle de Versão Distribuído.

Principais recursos:
• Commits
• Branches
• Merge
• Histórico completo

Criado por Linus Torvalds em 2005."""
)

# Slide 5

add_slide(
"Fluxo Básico do Git",
"""1. Desenvolvedor altera código
2. Realiza Commit
3. Envia para o repositório (Push)
4. Equipe compartilha mudanças

Resultado:
Maior organização e segurança."""
)

# Slide 6

add_slide(
"Versionamento de APIs",
"""APIs conectam sistemas diferentes.

Exemplo:
v1 → /api/v1/alunos
v2 → /api/v2/alunos

Benefícios:
• Compatibilidade
• Evolução gradual
• Menor impacto nos usuários"""
)

# Slide 7

add_slide(
"Aplicações no Mundo Real",
"""Empresas que utilizam versionamento:

• Google
• Microsoft
• Netflix
• Amazon
• Spotify

Todos os sistemas modernos utilizam controle de versões."""
)

# Slide 8

add_slide(
"Atividade",
"""1. O que aconteceria sem versionamento?

2. Qual a função de um commit?

3. Por que APIs possuem versões?

4. Cite uma vantagem do Git para equipes.

5. Como o versionamento contribui para a qualidade do software?"""
   )

prs.save("Aula_09_Versionamento_Codigo_APIs.pptx")
print("Apresentação criada com sucesso!")

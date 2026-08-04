from pptx import Presentation
from pptx.util import Pt

# Criando apresentação
prs = Presentation()

# Função para adicionar slides de conteúdo
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Título e conteúdo
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    title_box.text = title

    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.text = content

    # Ajustando tamanho da fonte
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Função para adicionar slides de atividade
def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    title_box.text = title

    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.text = questions

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)

# Slide 1 - Capa
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Instalação da Memória RAM"
slide.placeholders[1].text = "Montagem e Manutenção de Computadores"

# Slide 2 - O que é Memória RAM
add_slide(
    "O que é Memória RAM?",
    "A Memória RAM (Random Access Memory) é responsável pelo armazenamento temporário "
    "dos dados utilizados pelo computador.\n\n"
    "Ela permite que programas e o sistema operacional funcionem com rapidez e eficiência.\n\n"
    "Quanto maior a quantidade de RAM, maior a capacidade de executar várias tarefas simultaneamente."
)

# Slide 3 - Tipos de Memória RAM
add_slide(
    "Tipos de Memória RAM",
    "Existem diferentes tipos de memória RAM utilizados nos computadores:\n\n"
    "- DDR\n"
    "- DDR2\n"
    "- DDR3\n"
    "- DDR4\n"
    "- DDR5\n\n"
    "Cada tipo possui velocidades e encaixes diferentes."
)

# Slide 4 - Identificando o Slot
add_slide(
    "Identificando os Slots de Memória",
    "Os módulos de memória RAM são instalados nos slots da placa-mãe.\n\n"
    "Os slots geralmente ficam próximos ao processador.\n\n"
    "É importante verificar:\n"
    "- O tipo de memória compatível\n"
    "- A capacidade suportada\n"
    "- A posição correta do encaixe"
)

# Slide 5 - Cuidados Antes da Instalação
add_slide(
    "Cuidados Antes da Instalação",
    "Antes de instalar a memória RAM devemos:\n\n"
    "- Desligar o computador\n"
    "- Retirar o cabo de energia\n"
    "- Utilizar pulseira antiestática, se possível\n"
    "- Evitar tocar nos contatos dourados da memória"
)

# Slide 6 - Instalando a Memória RAM
add_slide(
    "Instalação da Memória RAM",
    "Para instalar a memória RAM:\n\n"
    "1. Abra as travas laterais do slot\n"
    "2. Alinhe o encaixe da memória com o slot\n"
    "3. Pressione cuidadosamente até ouvir o clique das travas\n\n"
    "A memória deve ficar totalmente encaixada."
)

# Slide 7 - Verificação da Instalação
add_slide(
    "Verificando a Instalação",
    "Após instalar a memória RAM:\n\n"
    "- Ligue o computador\n"
    "- Verifique se o sistema reconheceu a memória\n"
    "- Confira as informações na BIOS ou no sistema operacional\n\n"
    "Caso não reconheça, revise o encaixe do módulo."
)

# Slide 8 - Problemas Comuns
add_slide(
    "Problemas Comuns",
    "Alguns problemas podem ocorrer durante a instalação:\n\n"
    "- Memória mal encaixada\n"
    "- Incompatibilidade com a placa-mãe\n"
    "- Sujeira nos contatos\n"
    "- Defeito no módulo RAM\n\n"
    "Esses problemas podem impedir o funcionamento do computador."
)

# Slide 09 - Encerramento
add_slide(
    "Conclusão",
    "A instalação correta da memória RAM é fundamental para o desempenho e estabilidade do computador.\n\n"
    "Seguir os procedimentos adequados evita danos aos componentes e garante o funcionamento correto do sistema."
)

# Salvando apresentação
file_path = "aula_10_instalacao_memoria_ram.pptx"
prs.save(file_path)

print(f"Apresentação salva com sucesso em: {file_path}")
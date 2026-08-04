from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Slide 1
add_slide(
    "Autenticação em APIs com Flask",
    "Exemplo prático de autenticação usando Flask e sessões.\n\n"
    "Tópicos:\n"
    "- Autenticação e autorização\n"
    "- Implementação com Flask\n"
    "- OAuth, JWT e API Keys\n"
    "- Depuração\n"
    "- Documentação com Swagger"
)

# Slide 2
add_slide(
    "Autenticação x Autorização",
    "Autenticação: verifica quem é o usuário.\n\n"
    "Exemplo: login com usuário e senha.\n\n"
    "Autorização: define o que o usuário pode acessar após autenticado.\n\n"
    "Exemplo: administrador pode acessar recursos que usuários comuns não podem."
)

# Slide 3
add_slide(
    "Visão Geral do Código Flask",
    "A aplicação utiliza:\n"
    "- Flask para criar o servidor web.\n"
    "- Session para armazenar o estado do login.\n"
    "- Rotas para login, logout e página principal.\n"
    "- Redirecionamentos para controlar o acesso."
)

# Slide 4
add_slide(
    "Configuração Inicial",
    "app = Flask(__name__)\n"
    "app.secret_key = 'segredo'\n\n"
    "A secret_key protege os dados armazenados na sessão.\n"
    "Sem ela o Flask não consegue assinar os cookies de sessão."
)

# Slide 5
add_slide(
    "Credenciais de Exemplo",
    "USUARIO = 'admin'\n"
    "SENHA = '123456'\n\n"
    "Neste exemplo as credenciais são armazenadas diretamente no código.\n\n"
    "Em sistemas reais elas devem ficar em um banco de dados com senhas criptografadas."
)

# Slide 6
add_slide(
    "Rota Principal (/)",
    "Verifica se existe 'usuario' na sessão.\n\n"
    "Se existir:\n"
    "- Exibe mensagem de boas-vindas.\n"
    "- Permite logout.\n\n"
    "Caso contrário:\n"
    "- Redireciona para a tela de login."
)

# Slide 7
add_slide(
    "Processo de Login",
    "A rota /login aceita GET e POST.\n\n"
    "GET:\n"
    "- Exibe formulário de login.\n\n"
    "POST:\n"
    "- Recebe usuário e senha.\n"
    "- Compara com as credenciais cadastradas.\n"
    "- Cria a sessão do usuário."
)

# Slide 8
add_slide(
    "Uso da Sessão",
    "session['usuario'] = usuario\n\n"
    "A sessão mantém o usuário autenticado entre diferentes requisições.\n\n"
    "Enquanto a sessão existir, o usuário continuará logado."
)

# Slide 9
add_slide(
    "Logout",
    "session.pop('usuario', None)\n\n"
    "Remove o usuário da sessão.\n\n"
    "Após isso o acesso às páginas protegidas exige novo login."
)

# Slide 10
add_slide(
    "Fluxo da Autenticação",
    "1. Usuário acessa o sistema.\n"
    "2. É redirecionado para login.\n"
    "3. Informa usuário e senha.\n"
    "4. Sistema valida credenciais.\n"
    "5. Sessão é criada.\n"
    "6. Usuário acessa recursos protegidos."
)

# Slide 11
add_slide(
    "OAuth 2.0",
    "OAuth permite autenticação e autorização sem expor credenciais.\n\n"
    "Utiliza tokens de acesso.\n\n"
    "Exemplos:\n"
    "- Login com Google\n"
    "- Login com GitHub\n"
    "- Login com Facebook"
)

# Slide 12
add_slide(
    "JWT (JSON Web Token)",
    "JWT é um padrão para troca segura de informações.\n\n"
    "Fluxo:\n"
    "- Usuário faz login.\n"
    "- Servidor gera um token.\n"
    "- Cliente envia o token em cada requisição.\n"
    "- Servidor valida o token."
)

# Slide 13
add_slide(
    "Vantagens do JWT",
    "- Não exige armazenamento de sessão no servidor.\n"
    "- Escalabilidade maior.\n"
    "- Muito utilizado em APIs REST.\n"
    "- Fácil integração entre sistemas."
)

# Slide 14
add_slide(
    "API Keys",
    "Método simples de autenticação.\n\n"
    "Exemplo:\n"
    "GET /produtos?api_key=12345678\n\n"
    "Desvantagem:\n"
    "- A chave pode ficar exposta.\n"
    "- Menor nível de segurança."
)

# Slide 15
add_slide(
    "Comparação dos Métodos",
    "Sessão Flask:\n"
    "- Simples\n"
    "- Boa para aplicações web\n\n"
    "JWT:\n"
    "- Ideal para APIs REST\n"
    "- Stateless\n\n"
    "OAuth:\n"
    "- Integração com provedores externos\n\n"
    "API Keys:\n"
    "- Simples, porém menos segura"
)

# Slide 16
add_slide(
    "Técnicas de Depuração",
    "Debug é o processo de localizar e corrigir erros.\n\n"
    "Principais técnicas:\n"
    "- Logs\n"
    "- Breakpoints\n"
    "- Testes automatizados\n"
    "- Monitoramento em tempo real"
)

# Slide 17
add_slide(
    "Modo Debug do Flask",
    "app.run(debug=True)\n\n"
    "Benefícios:\n"
    "- Recarregamento automático.\n"
    "- Exibição detalhada de erros.\n"
    "- Facilita o desenvolvimento."
)

# Slide 18
add_slide(
    "Testes Automatizados",
    "Ferramentas como Jest, PyTest e Unittest permitem:\n\n"
    "- Detectar erros rapidamente.\n"
    "- Validar funcionalidades.\n"
    "- Aumentar a qualidade do sistema."
)

# Slide 19
add_slide(
    "Documentação de APIs",
    "Documentar uma API facilita:\n\n"
    "- Uso por outros desenvolvedores.\n"
    "- Integrações.\n"
    "- Manutenção.\n"
    "- Testes."
)

# Slide 20
add_slide(
    "OpenAPI Specification (OAS)",
    "Padrão para descrição de APIs REST.\n\n"
    "Permite documentação legível por humanos e sistemas.\n\n"
    "Base para ferramentas modernas de documentação."
)

# Slide 21
add_slide(
    "Swagger",
    "Ferramenta que implementa OpenAPI.\n\n"
    "Recursos:\n"
    "- Documentação interativa.\n"
    "- Teste de endpoints.\n"
    "- Geração automática de documentação."
)

# Slide 22
add_slide(
    "Conclusão",
    "O exemplo Flask demonstra autenticação baseada em sessão.\n\n"
    "Para sistemas maiores, recomenda-se:\n"
    "- Senhas criptografadas.\n"
    "- JWT para APIs REST.\n"
    "- OAuth para integração externa.\n"
    "- Swagger para documentação."
)

file_path = "aula_06_autenticacao_flask_api.pptx"
prs.save(file_path)

print(file_path)
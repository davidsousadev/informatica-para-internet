from pptx import Presentation

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = questions

# Slide 1
add_slide(
    "Integração Front-end e Back-end",
    "Uma aplicação web é composta por duas partes principais:\n\n"
    "• Front-end: interface utilizada pelo usuário.\n"
    "• Back-end: responsável pelo processamento das informações.\n\n"
    "A comunicação entre essas camadas permite enviar dados, processar requisições e exibir resultados."
)

# Slide 2
add_slide(
    "Arquiteturas Web",
    "As aplicações podem ser:\n\n"
    "• Monolíticas: front-end e back-end fazem parte da mesma aplicação.\n"
    "• Baseadas em microsserviços: a comunicação ocorre por APIs.\n\n"
    "Frameworks Python como Flask permitem criar aplicações web de forma simples, seguindo a arquitetura MVC."
)

# Slide 3
add_slide(
    "Exemplo de Integração em Python Puro",
    "Front-end (HTML):\n"
    "• Possui um botão.\n"
    "• Envia uma requisição usando fetch().\n\n"
    "Back-end (Python):\n"
    "• Recebe a requisição.\n"
    "• Retorna a mensagem 'Mensagem do back-end Python'.\n\n"
    "Fluxo:\n"
    "Usuário → Front-end → Back-end → Resposta → Front-end"
)

# Slide 4
add_activity_slide(
    "Atividade",
    "1. O que é front-end?\n\n"
    "2. O que é back-end?\n\n"
    "3. Qual a diferença entre arquitetura monolítica e microsserviços?\n\n"
    "4. No exemplo apresentado, qual tecnologia foi utilizada para o back-end?\n\n"
    "5. O que acontece quando o usuário clica no botão?"
)

file_path = "aula_08_integracao_frontend_backend.pptx"
prs.save(file_path)

print(file_path)
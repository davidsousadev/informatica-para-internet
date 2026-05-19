from pptx import Presentation
from pptx.util import Inches, Pt

# Criando apresentação
prs = Presentation()

# Função para adicionar slides de conteúdo
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Título e conteúdo
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    title_box.text = title

    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.text = content

    # Ajustando tamanho da fonte
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Função para adicionar slides de atividade
def add_activity_slide(title, questions):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.title
    title_box.text = title

    content_box = slide.placeholders[1]
    text_frame = content_box.text_frame
    text_frame.text = questions

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)

# Slide 1 - Capa
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Conexão de Cabos e Configuração de Drives IDE"
slide.placeholders[1].text = "Montagem e Manutenção de Computadores"

# Slide 2 - Painel Frontal
add_slide(
    "Conexão do Painel Frontal",
    "Após fixar a placa-mãe, conectamos os cabos do painel frontal do gabinete.\n\n"
    "Alguns gabinetes possuem conectores unificados, porém a maioria utiliza conectores separados.\n\n"
    "Cada cabo possui identificação própria, como:\n"
    "- Power LED\n"
    "- Reset\n"
    "- HDD LED\n\n"
    "A conexão correta deve seguir a ordem dos pinos da placa-mãe."
)

# Slide 3 - Manual e USB
add_slide(
    "Identificação dos Pinos",
    "A posição correta dos pinos pode ser encontrada:\n\n"
    "- Na própria placa-mãe\n"
    "- No manual da placa-mãe\n\n"
    "Após isso, devemos localizar os conectores USB do gabinete e conectá-los corretamente à placa-mãe."
)

# Slide 4 - Drives IDE
add_slide(
    "Conexão de Drives IDE",
    "Após conectar os cabos do painel frontal, conectamos os dispositivos de armazenamento:\n\n"
    "- HD\n"
    "- CD\n"
    "- DVD\n\n"
    "Nos drives IDE é necessário configurar os Jumpers localizados na parte traseira do dispositivo."
)

# Slide 5 - Interfaces IDE
add_slide(
    "Interfaces IDE",
    "A placa-mãe pode possuir 1 ou 2 interfaces IDE:\n\n"
    "- IDE 0 (Primária)\n"
    "- IDE 1 (Secundária)\n\n"
    "Cada interface permite conectar até 2 dispositivos IDE."
)

# Slide 6 - Master e Slave
add_slide(
    "Configuração Master e Slave",
    "Os dispositivos IDE devem ser configurados como:\n\n"
    "- Master (Mestre)\n"
    "- Slave (Escravo)\n\n"
    "Se houver dois dispositivos na mesma interface:\n"
    "- Um deve ser Master\n"
    "- O outro deve ser Slave\n\n"
    "Se houver apenas um dispositivo, ele deve ficar configurado como Master."
)

# Slide 7 - Configuração dos Jumpers
add_slide(
    "Configuração dos Jumpers",
    "A configuração dos Jumpers varia conforme:\n\n"
    "- Marca do HD\n"
    "- Modelo do dispositivo\n\n"
    "As instruções normalmente podem ser encontradas:\n"
    "- Na etiqueta do dispositivo\n"
    "- No manual do fabricante"
)

# Slide 8 - Atividade
add_activity_slide(
    "Atividade",
    "1. Qual a função dos cabos do painel frontal?\n\n"
    "2. Onde encontramos a ordem correta dos pinos da placa-mãe?\n\n"
    "3. O que significa Master e Slave em dispositivos IDE?\n\n"
    "4. Como deve ser configurado um único dispositivo IDE?"
)

# Slide 9 - Encerramento
add_slide(
    "Conclusão",
    "A correta conexão dos cabos e configuração dos dispositivos IDE "
    "é essencial para o funcionamento adequado do computador.\n\n"
    "A atenção aos detalhes durante a montagem evita falhas e melhora a manutenção do sistema."
)

# Salvando apresentação
file_path = "conexao_painel_frontal_e_drives_IDE.pptx"
prs.save(file_path)

print(f"Apresentação salva com sucesso em: {file_path}")
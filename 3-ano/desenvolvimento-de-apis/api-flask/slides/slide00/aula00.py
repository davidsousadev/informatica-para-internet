from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content_list):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, item in enumerate(content_list):
        if i == 0:
            text_frame.text = item
        else:
            p = text_frame.add_paragraph()
            p.text = item
            p.level = 0

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

# Slide inicial
add_title_slide(
    "Entendendo um código Flask",
    "Explicação linha por linha para Ensino Médio"
)

# Slides
add_slide("O que é Flask?", [
    "Flask é um framework para criar sites usando Python",
    "Framework = conjunto de ferramentas prontas",
    "Ajuda a criar aplicações web de forma mais simples",
    "[ESPAÇO PARA LOGO DO FLASK]"
])

add_slide("Importando o Flask", [
    "Código: from flask import Flask",
    "from = pegar algo de uma biblioteca",
    "flask = biblioteca usada para criar sites",
    "Flask = ferramenta principal da biblioteca"
])

add_slide("Criando a aplicação", [
    "Código: app = Flask(__name__)",
    "app é o nome da aplicação",
    "Flask(__name__) cria o projeto web",
    "__name__ ajuda o Python a identificar o arquivo"
])

add_slide("Decorador de rota", [
    "Código: @app.route('/')",
    "Define o endereço da página",
    "/ significa página inicial do site",
    "Quando alguém acessar o site, esta função será executada"
])

add_slide("Criando uma função", [
    "Código: def hello_world():",
    "def é usado para criar funções",
    "Função = bloco de código reutilizável",
    "hello_world é o nome da função"
])

add_slide("Indentação no Python", [
    "Tudo abaixo da função possui espaço no começo da linha",
    "Isso se chama indentação",
    "Python usa indentação para organizar blocos de código",
    "[ESPAÇO PARA EXEMPLO VISUAL]"
])

add_slide("Retornando uma mensagem", [
    "Código: return '<p>Hello, World!</p>'",
    "return envia uma resposta para o navegador",
    "<p> é uma tag HTML de parágrafo",
    "O navegador mostrará: Hello, World!"
])

add_slide("HTML dentro do Python", [
    "O texto retornado está em HTML",
    "HTML é a linguagem usada para criar páginas web",
    "Python envia o HTML para o navegador",
    "[ESPAÇO PARA MOSTRAR HTML SIMPLES]"
])

add_slide("Executando o servidor", [
    "Código: app.run()",
    "Inicia o servidor Flask",
    "Servidor = programa que responde acessos ao site",
    "Depois disso, o site pode ser aberto no navegador"
])

add_slide("Fluxo do Programa", [
    "1. Flask é importado",
    "2. A aplicação é criada",
    "3. A rota é definida",
    "4. O servidor inicia"
])

add_slide("O que acontece ao acessar o site?", [
    "Usuário abre o navegador",
    "Acessa o endereço do servidor Flask",
    "Flask executa hello_world()",
    "O navegador recebe a mensagem Hello, World!"
])

add_slide("Resumo Geral", [
    "Flask permite criar sites com Python",
    "Rotas definem páginas",
    "Funções executam ações",
    "HTML é enviado ao navegador"
])

add_slide("Código Completo", [
    "from flask import Flask",
    "app = Flask(__name__)",
    "@app.route('/')",
    "def hello_world(): return '<p>Hello, World!</p>'",
    "app.run()"
])

# Salvar
file_path = "explicacao-flask.pptx"
prs.save(file_path)

print("Arquivo criado:", file_path)